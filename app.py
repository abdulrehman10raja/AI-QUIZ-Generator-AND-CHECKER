import streamlit as st
import fitz
from groq import Groq
from dotenv import load_dotenv
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
import os, json, re, io, base64, time
import google.auth.transport.requests
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from pptx import Presentation
from docx import Document as DocxDocument

load_dotenv()

st.set_page_config(
    page_title="QUEST // AI Quiz System",
    page_icon="⬡",
    layout="wide",
    initial_sidebar_state="expanded"
)

if "page" not in st.session_state:
    st.session_state.page = "splash"

current_page = st.session_state.page

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Share+Tech+Mono&family=JetBrains+Mono:wght@300;400;500;600;700&display=swap');

:root {
    --bg-base:      #060D0F;
    --bg-surface:   #0C1A1E;
    --bg-raised:    #0F2128;
    --bg-hover:     #142830;
    --cyan:         #00E5FF;
    --cyan-dim:     #009AB0;
    --cyan-glow:    rgba(0, 229, 255, 0.12);
    --cyan-border:  rgba(0, 229, 255, 0.25);
    --cyan-strong:  rgba(0, 229, 255, 0.5);
    --purple:       #7C3AFF;
    --purple-dim:   rgba(124, 58, 255, 0.2);
    --green:        #00FF88;
    --green-dim:    rgba(0, 255, 136, 0.15);
    --red:          #FF4444;
    --red-dim:      rgba(255, 68, 68, 0.15);
    --amber:        #FFCB3D;
    --ink:          #C8F0F7;
    --ink-bright:   #E8F8FF;
    --ink-muted:    #3D6B78;
    --ink-faint:    #1D3A42;
}

* { scroll-behavior: smooth; box-sizing: border-box; }

.stApp {
    background-color: var(--bg-base) !important;
    background-image:
        linear-gradient(rgba(0,229,255,0.03) 1px, transparent 1px),
        linear-gradient(90deg, rgba(0,229,255,0.03) 1px, transparent 1px) !important;
    background-size: 40px 40px !important;
    font-family: 'JetBrains Mono', monospace !important;
    color: var(--ink) !important;
}

.main .block-container { padding-top: 1rem; max-width: 1100px; }

::-webkit-scrollbar { width: 6px; }
::-webkit-scrollbar-track { background: var(--bg-base); }
::-webkit-scrollbar-thumb { background: var(--cyan-dim); border-radius: 3px; }

@keyframes pulse-cyan {
    0%, 100% { box-shadow: 0 0 8px var(--cyan-glow), inset 0 0 8px var(--cyan-glow); }
    50%       { box-shadow: 0 0 22px rgba(0,229,255,0.28), inset 0 0 14px rgba(0,229,255,0.1); }
}
@keyframes scan-line {
    0%   { transform: translateY(-100%); opacity: 0; }
    15%  { opacity: 0.4; }
    85%  { opacity: 0.4; }
    100% { transform: translateY(1000%); opacity: 0; }
}
@keyframes blink-cursor {
    0%, 100% { opacity: 1; }
    50%       { opacity: 0; }
}
@keyframes flicker-in {
    0%   { opacity: 0; }
    20%  { opacity: 0.8; }
    40%  { opacity: 0.3; }
    60%  { opacity: 1; }
    80%  { opacity: 0.7; }
    100% { opacity: 1; }
}
@keyframes slide-up {
    from { opacity: 0; transform: translateY(16px); }
    to   { opacity: 1; transform: translateY(0); }
}
@keyframes slide-right {
    from { opacity: 0; transform: translateX(-12px); }
    to   { opacity: 1; transform: translateX(0); }
}
@keyframes glow-pulse {
    0%, 100% { text-shadow: 0 0 8px rgba(0,229,255,0.6); }
    50%       { text-shadow: 0 0 24px rgba(0,229,255,0.9), 0 0 48px rgba(0,229,255,0.3); }
}
@keyframes border-flow {
    0%   { background-position: 0% 50%; }
    50%  { background-position: 100% 50%; }
    100% { background-position: 0% 50%; }
}
@keyframes dot-blink {
    0%, 100% { opacity: 1; }
    33%       { opacity: 0.2; }
    66%       { opacity: 0.7; }
}
@keyframes rotate-slow {
    from { transform: rotate(0deg); }
    to   { transform: rotate(360deg); }
}
@keyframes status-ping {
    0%   { box-shadow: 0 0 0 0 rgba(0,255,136,0.7); }
    70%  { box-shadow: 0 0 0 8px rgba(0,255,136,0); }
    100% { box-shadow: 0 0 0 0 rgba(0,255,136,0); }
}
@keyframes mega-glow {
    0%, 100% {
        text-shadow: 0 0 10px rgba(0,229,255,0.6), 0 0 20px rgba(0,229,255,0.4), 0 0 40px rgba(0,229,255,0.2);
    }
    50% {
        text-shadow: 0 0 20px rgba(0,229,255,1), 0 0 40px rgba(0,229,255,0.7), 0 0 80px rgba(0,229,255,0.4), 0 0 120px rgba(124,58,255,0.3);
    }
}
@keyframes border-rotate {
    0%   { background-position: 0% 0%; }
    100% { background-position: 200% 0%; }
}
@keyframes float-y {
    0%, 100% { transform: translateY(0); }
    50%       { transform: translateY(-8px); }
}
@keyframes float-y2 {
    0%, 100% { transform: translateY(0); }
    50%       { transform: translateY(-5px); }
}
@keyframes badge-shimmer {
    0%   { background-position: -200% 0; }
    100% { background-position: 200% 0; }
}
@keyframes ring-spin {
    from { transform: rotate(0deg); }
    to   { transform: rotate(360deg); }
}
@keyframes ring-spin-reverse {
    from { transform: rotate(360deg); }
    to   { transform: rotate(0deg); }
}
@keyframes avatar-pulse {
    0%, 100% { box-shadow: 0 0 0 0 rgba(0,229,255,0.5), 0 0 20px rgba(124,58,255,0.4); }
    70%       { box-shadow: 0 0 0 10px rgba(0,229,255,0), 0 0 30px rgba(124,58,255,0.6); }
}
@keyframes tag-glow-cycle {
    0%, 100% { box-shadow: 0 0 6px currentColor; }
    50%       { box-shadow: 0 0 16px currentColor; }
}
@keyframes flow-card-lift {
    0%, 100% { transform: translateY(0); box-shadow: 0 0 0 rgba(0,229,255,0); }
    33%       { transform: translateY(-8px); box-shadow: 0 12px 28px rgba(0,229,255,0.12); }
}
@keyframes stat-count {
    0%   { opacity: 0.3; transform: scale(0.9); }
    100% { opacity: 1; transform: scale(1); }
}
@keyframes pill-glow {
    0%, 100% { border-color: var(--ink-faint); color: var(--ink-muted); box-shadow: none; }
    50%       { border-color: var(--cyan-dim); color: var(--cyan-dim); box-shadow: 0 0 12px rgba(0,229,255,0.15); }
}
@keyframes sb-item-slide {
    from { opacity: 0; transform: translateX(-16px); }
    to   { opacity: 1; transform: translateX(0); }
}
@keyframes sb-glow-pulse {
    0%, 100% { box-shadow: 0 0 0 rgba(0,229,255,0); }
    50%       { box-shadow: 0 0 18px rgba(0,229,255,0.15); }
}
@keyframes sidebar-scan {
    0%   { top: -10%; opacity: 0; }
    10%  { opacity: 0.6; }
    90%  { opacity: 0.6; }
    100% { top: 110%; opacity: 0; }
}
@keyframes logo-flicker {
    0%, 92%, 100% { opacity: 1; }
    93%           { opacity: 0.4; }
    95%           { opacity: 1; }
    97%           { opacity: 0.2; }
    99%           { opacity: 1; }
}
@keyframes creator-float {
    0%, 100% { transform: translateY(0) rotate(0deg); }
    25%       { transform: translateY(-4px) rotate(1deg); }
    75%       { transform: translateY(4px) rotate(-1deg); }
}
@keyframes hud-rotate {
    0%   { opacity: 0.3; }
    50%  { opacity: 0.8; }
    100% { opacity: 0.3; }
}
@keyframes progress-bar-anim {
    0%   { background-position: 0 0; }
    100% { background-position: 40px 0; }
}
@keyframes typing {
    from { width: 0; }
    to   { width: 100%; }
}

@keyframes splash-hex-spin {
    0%   { transform: rotate(0deg) scale(1); }
    25%  { transform: rotate(90deg) scale(1.08); }
    50%  { transform: rotate(180deg) scale(1); }
    75%  { transform: rotate(270deg) scale(1.05); }
    100% { transform: rotate(360deg) scale(1); }
}
@keyframes splash-ring-1 {
    0%   { transform: rotate(0deg);   opacity: 0.6; }
    100% { transform: rotate(360deg); opacity: 0.6; }
}
@keyframes splash-ring-2 {
    0%   { transform: rotate(0deg);   opacity: 0.35; }
    100% { transform: rotate(-360deg);opacity: 0.35; }
}
@keyframes splash-ring-3 {
    0%   { transform: rotate(0deg);   opacity: 0.2; }
    100% { transform: rotate(360deg); opacity: 0.2; }
}
@keyframes splash-float {
    0%, 100% { transform: translateY(0px); }
    50%       { transform: translateY(-14px); }
}
@keyframes splash-title-reveal {
    0%   { opacity: 0; letter-spacing: 14px; filter: blur(6px); }
    100% { opacity: 1; letter-spacing: 3px;  filter: blur(0); }
}
@keyframes splash-sub-reveal {
    0%   { opacity: 0; transform: translateY(10px); }
    100% { opacity: 1; transform: translateY(0); }
}
@keyframes splash-bar-slide {
    from { width: 0; opacity: 0; }
    to   { width: 120px; opacity: 1; }
}
@keyframes splash-border-run {
    0%   { background-position: 0% 0%; }
    100% { background-position: 300% 0%; }
}
@keyframes splash-creator-in {
    0%   { opacity: 0; transform: scale(0.92) translateY(8px); }
    100% { opacity: 1; transform: scale(1) translateY(0); }
}
@keyframes splash-scan {
    0%   { transform: translateY(-100%); opacity: 0; }
    15%  { opacity: 0.5; }
    85%  { opacity: 0.5; }
    100% { transform: translateY(2000%); opacity: 0; }
}
@keyframes splash-btn-pulse {
    0%, 100% {
        box-shadow: 0 0 16px rgba(0,229,255,0.25), 0 0 0 0 rgba(0,229,255,0.4);
    }
    50% {
        box-shadow: 0 0 32px rgba(0,229,255,0.5), 0 0 0 8px rgba(0,229,255,0);
    }
}
@keyframes splash-particle {
    0%   { opacity: 0; transform: translateY(0) scale(0); }
    20%  { opacity: 1; }
    80%  { opacity: 0.6; }
    100% { opacity: 0; transform: translateY(-80px) scale(1.5); }
}
@keyframes splash-grid-pulse {
    0%, 100% { opacity: 0.025; }
    50%       { opacity: 0.055; }
}
@keyframes splash-corner-blink {
    0%, 100% { opacity: 0.9; }
    50%       { opacity: 0.3; }
}
@keyframes splash-tag-in {
    0%   { opacity: 0; transform: scale(0.8); }
    60%  { transform: scale(1.06); }
    100% { opacity: 1; transform: scale(1); }
}
@keyframes splash-streak {
    0%   { left: -30%; opacity: 0; }
    10%  { opacity: 0.7; }
    90%  { opacity: 0.7; }
    100% { left: 130%; opacity: 0; }
}

.splash-page [data-testid="stSidebar"] { display: none !important; }

.splash-hide-sidebar [data-testid="stSidebar"],
.splash-hide-sidebar [data-testid="stSidebarCollapsedControl"],
.splash-hide-sidebar [data-testid="stSidebarCollapseButton"] {
    display: none !important;
}

[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #040A0C 0%, #060F14 40%, #04090B 100%) !important;
    border-right: 1px solid var(--cyan-border) !important;
    box-shadow: 4px 0 40px rgba(0,229,255,0.08) !important;
    position: relative;
    overflow: hidden;
}
[data-testid="stSidebar"]::before {
    content: "";
    position: fixed;
    top: 0; left: 0;
    width: 280px; height: 100vh;
    background-image:
        linear-gradient(rgba(0,229,255,0.025) 1px, transparent 1px),
        linear-gradient(90deg, rgba(0,229,255,0.025) 1px, transparent 1px);
    background-size: 18px 18px;
    pointer-events: none;
    z-index: 0;
}
[data-testid="stSidebar"]::after {
    content: "";
    position: fixed;
    left: 0;
    width: 280px;
    height: 3px;
    background: linear-gradient(90deg, transparent, rgba(0,229,255,0.5), transparent);
    animation: sidebar-scan 5s ease-in-out infinite;
    pointer-events: none;
    z-index: 2;
}
[data-testid="stSidebar"] > div { position: relative; z-index: 1; }

[data-testid="stSidebar"] .stMarkdown p,
[data-testid="stSidebar"] .stMarkdown li {
    color: var(--ink) !important;
    font-family: 'JetBrains Mono', monospace !important;
}

[data-testid="stSidebarCollapsedControl"],
[data-testid="stSidebarCollapseButton"] {
    position: fixed !important; top: 16px !important; left: 12px !important;
    z-index: 999999 !important; background: var(--bg-surface) !important;
    border: 1px solid var(--cyan-border) !important; border-radius: 6px !important;
    width: 42px !important; height: 42px !important;
    box-shadow: 0 0 14px var(--cyan-glow) !important;
}
[data-testid="stSidebarCollapsedControl"] button,
[data-testid="stSidebarCollapseButton"] button {
    background: transparent !important; border: none !important; box-shadow: none !important;
    width: 42px !important; height: 42px !important; padding: 0 !important;
    color: transparent !important; font-size: 0 !important;
    display: flex !important; align-items: center !important; justify-content: center !important;
    position: relative !important;
}
[data-testid="stSidebarCollapsedControl"] button svg,
[data-testid="stSidebarCollapseButton"] button svg { display: none !important; }
[data-testid="stSidebarCollapsedControl"] button::before,
[data-testid="stSidebarCollapseButton"] button::before {
    content: "" !important; display: block !important; width: 16px !important; height: 2px !important;
    background: var(--cyan) !important; border-radius: 1px !important;
    box-shadow: 0 5px 0 var(--cyan), 0 10px 0 var(--cyan) !important;
    position: absolute !important; top: 50% !important; left: 50% !important;
    transform: translate(-50%, -5px) !important;
}

.sidebar-logo {
    font-family: 'Share Tech Mono', monospace;
    font-size: 1.6rem;
    color: var(--cyan);
    text-align: center;
    padding: 1.2rem 0 0.2rem 0;
    letter-spacing: 4px;
    text-shadow: 0 0 20px rgba(0,229,255,0.7), 0 0 40px rgba(0,229,255,0.3);
    animation: logo-flicker 6s ease-in-out infinite, glow-pulse 3s ease-in-out infinite;
    cursor: default;
}
.sidebar-sub {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.6rem;
    color: var(--ink-muted);
    text-align: center;
    letter-spacing: 3px;
    margin-bottom: 0.5rem;
    text-transform: uppercase;
    animation: sb-item-slide 0.6s ease-out 0.1s backwards;
}

.sidebar-status {
    display: flex; align-items: center; gap: 8px;
    padding: 0.4rem 0.8rem;
    margin: 0.5rem 0.5rem 0 0.5rem;
    background: rgba(0,255,136,0.05);
    border: 1px solid rgba(0,255,136,0.15);
    border-radius: 4px;
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.65rem; color: var(--green); letter-spacing: 1px;
    animation: sb-item-slide 0.6s ease-out 0.2s backwards, sb-glow-pulse 3s ease-in-out infinite 1s;
}
.status-dot {
    width: 6px; height: 6px; border-radius: 50%;
    background: var(--green);
    animation: status-ping 1.8s ease-in-out infinite;
    flex-shrink: 0;
}

.sidebar-divider {
    border: none; border-top: 1px solid var(--cyan-border); margin: 0.8rem 0;
    position: relative;
}
.sidebar-divider::after {
    content: "◆"; position: absolute; left: 50%; top: -7px;
    transform: translateX(-50%); background: #060F14;
    color: var(--cyan-dim); font-size: 0.55rem; padding: 0 6px;
    animation: dot-blink 3s ease-in-out infinite;
}

.quest-tracker { display: flex; align-items: center; justify-content: space-between; padding: 0.6rem 0.8rem 1rem 0.8rem; }
.quest-node { display: flex; flex-direction: column; align-items: center; gap: 5px; flex: 1; }
.quest-circle {
    width: 32px; height: 32px; border-radius: 4px;
    display: flex; align-items: center; justify-content: center;
    font-family: 'Share Tech Mono', monospace; font-size: 0.85rem;
    border: 1px solid var(--ink-muted); background: var(--bg-surface);
    color: var(--ink-muted); transition: all .2s ease;
}
.quest-circle.done {
    background: rgba(0,255,136,0.1); border-color: var(--green); color: var(--green);
    box-shadow: 0 0 10px rgba(0,255,136,0.2);
}
.quest-circle.active {
    background: rgba(0,229,255,0.1); border-color: var(--cyan); color: var(--cyan);
    box-shadow: 0 0 14px rgba(0,229,255,0.3);
    animation: pulse-cyan 2s ease-in-out infinite;
}
.quest-line { flex: 0.6; height: 1px; background: var(--ink-faint); align-self: center; margin-top: -18px; }
.quest-line.done { background: var(--green); box-shadow: 0 0 6px rgba(0,255,136,0.3); }
.quest-label { font-family: 'JetBrains Mono', monospace; font-size: 0.52rem; color: var(--ink-muted); font-weight: 500; letter-spacing: 1px; text-transform: uppercase; }
.quest-label.active { color: var(--cyan); }

[data-testid="stSidebar"] .stButton > button {
    width: 100% !important; text-align: left !important;
    background: transparent !important; border: 1px solid var(--ink-faint) !important;
    border-radius: 4px !important; padding: 0.7rem 1rem !important; margin-bottom: 0.4rem !important;
    font-family: 'JetBrains Mono', monospace !important; font-weight: 500 !important;
    font-size: 0.78rem !important; color: var(--ink-muted) !important; cursor: pointer !important;
    transition: all 0.2s ease !important; box-shadow: none !important;
    letter-spacing: 0.5px !important; text-transform: uppercase !important;
    position: relative !important; overflow: hidden !important;
}
[data-testid="stSidebar"] .stButton > button::before {
    content: "" !important; position: absolute !important;
    top: 0 !important; left: -100% !important; width: 60% !important; height: 100% !important;
    background: linear-gradient(120deg, transparent, rgba(0,229,255,0.1), transparent) !important;
    transform: skewX(-20deg) !important;
    transition: left 0.5s ease !important;
}
[data-testid="stSidebar"] .stButton > button:hover::before { left: 150% !important; }
[data-testid="stSidebar"] .stButton > button:hover {
    background: var(--cyan-glow) !important;
    border-color: var(--cyan-border) !important;
    color: var(--cyan) !important;
    transform: translateX(4px) !important;
    box-shadow: inset 3px 0 0 var(--cyan), 0 0 16px rgba(0,229,255,0.08) !important;
}
[data-testid="stSidebar"] .stButton > button:focus { box-shadow: none !important; outline: none !important; }
[data-testid="stSidebar"] .nav-active > button {
    background: rgba(0,229,255,0.07) !important; border-color: var(--cyan) !important;
    color: var(--cyan) !important;
    box-shadow: inset 3px 0 0 var(--cyan), 0 0 14px rgba(0,229,255,0.08) !important;
    text-shadow: 0 0 8px rgba(0,229,255,0.5) !important;
    animation: sb-glow-pulse 2.5s ease-in-out infinite !important;
}
[data-testid="stSidebar"] .nav-active > button:hover { transform: none !important; }

.sidebar-creator {
    font-family: 'JetBrains Mono', monospace; font-size: 0.62rem;
    color: var(--ink-muted); text-align: center;
    padding: 0.8rem 0 0.5rem 0; letter-spacing: 1px;
    border-top: 1px solid var(--ink-faint); margin-top: 0.5rem;
    animation: sb-item-slide 0.6s ease-out 0.5s backwards;
}
.sidebar-creator span { color: var(--cyan-dim); }

.term-card {
    background: var(--bg-surface); border: 1px solid var(--cyan-border);
    border-top: 2px solid var(--cyan); border-radius: 6px;
    padding: 1.4rem 1.6rem; margin-bottom: 1.2rem; position: relative;
    animation: slide-up 0.35s ease-out; transition: border-color 0.2s ease, box-shadow 0.2s ease;
}
.term-card::before {
    content: ""; position: absolute; top: 0; left: 0; right: 0; height: 1px;
    background: linear-gradient(90deg, transparent, var(--cyan), transparent); opacity: 0.6;
}
.term-card:hover { border-color: var(--cyan-strong); box-shadow: 0 0 24px var(--cyan-glow), 0 0 1px var(--cyan-border); }
.term-card::after {
    content: ""; position: absolute; bottom: 6px; right: 6px;
    width: 10px; height: 10px;
    border-bottom: 1px solid var(--cyan-dim); border-right: 1px solid var(--cyan-dim); opacity: 0.5;
}

.sec-title {
    font-family: 'Share Tech Mono', monospace; font-size: 0.72rem;
    color: var(--cyan); letter-spacing: 3px; text-transform: uppercase;
    margin-bottom: 1rem; display: flex; align-items: center; gap: 8px;
}
.sec-title::before {
    content: "▸";
    color: var(--cyan-dim); font-size: 0.7rem;
}
.sec-title::after {
    content: ""; flex: 1; height: 1px;
    background: linear-gradient(90deg, var(--cyan-border), transparent);
}

.page-header { padding: 0.2rem 0 1.4rem 0; animation: flicker-in 0.6s ease-out; }
.page-title {
    font-family: 'Share Tech Mono', monospace; font-size: 1.8rem;
    color: var(--ink-bright); letter-spacing: 2px; margin-bottom: 0.3rem;
    text-shadow: 0 0 20px rgba(0,229,255,0.2);
}
.page-title span { color: var(--cyan); }
.page-subtitle { font-family: 'JetBrains Mono', monospace; font-size: 0.78rem; color: var(--ink-muted); letter-spacing: 1px; }

.journey-bar {
    display: flex; align-items: center; margin-bottom: 1.8rem;
    background: var(--bg-surface); border: 1px solid var(--cyan-border);
    border-radius: 6px; padding: 1rem 1.4rem; position: relative; overflow: hidden;
}
.journey-bar::after {
    content: ""; position: absolute; top: -50%; left: -50%;
    width: 20px; height: 200%;
    background: linear-gradient(90deg, transparent, rgba(0,229,255,0.06), transparent);
    animation: scan-line 6s linear infinite;
}
.journey-step { display: flex; align-items: center; gap: 0.7rem; flex: 1; animation: slide-right 0.4s ease-out backwards; }
.journey-step:nth-child(1) { animation-delay: 0.05s; }
.journey-step:nth-child(3) { animation-delay: 0.12s; }
.journey-step:nth-child(5) { animation-delay: 0.19s; }
.journey-circle {
    width: 38px; height: 38px; min-width: 38px; border-radius: 4px;
    display: flex; align-items: center; justify-content: center;
    font-family: 'Share Tech Mono', monospace; font-size: 1rem;
    border: 1px solid var(--ink-faint); background: var(--bg-raised); color: var(--ink-muted);
    transition: all .3s ease;
}
.journey-circle.done { background: rgba(0,255,136,0.08); border-color: var(--green); color: var(--green); box-shadow: 0 0 12px rgba(0,255,136,0.15); }
.journey-circle.active { background: rgba(0,229,255,0.08); border-color: var(--cyan); color: var(--cyan); box-shadow: 0 0 16px rgba(0,229,255,0.25); animation: pulse-cyan 2s ease-in-out infinite; }
.journey-text { display: flex; flex-direction: column; }
.journey-name { font-family: 'Share Tech Mono', monospace; font-size: 0.78rem; color: var(--ink); letter-spacing: 1px; text-transform: uppercase; }
.journey-sub { font-family: 'JetBrains Mono', monospace; font-size: 0.62rem; color: var(--ink-muted); letter-spacing: 0.5px; }
.journey-connector { flex: 0.4; height: 1px; background: var(--ink-faint); margin: 0 0.8rem; }
.journey-connector.done { background: var(--green); box-shadow: 0 0 6px rgba(0,255,136,0.3); }

.hero-wrapper {
    min-height: 88vh; display: flex; flex-direction: column;
    align-items: center; justify-content: center;
    padding: 2rem 1rem 3rem 1rem; text-align: center; position: relative;
}

.hud-corner {
    position: absolute; width: 60px; height: 60px;
    pointer-events: none; animation: hud-rotate 3s ease-in-out infinite;
}
.hud-corner.tl { top: 20px; left: 20px; border-top: 2px solid var(--cyan); border-left: 2px solid var(--cyan); animation-delay: 0s; }
.hud-corner.tr { top: 20px; right: 20px; border-top: 2px solid var(--cyan); border-right: 2px solid var(--cyan); animation-delay: 0.75s; }
.hud-corner.bl { bottom: 60px; left: 20px; border-bottom: 2px solid var(--cyan); border-left: 2px solid var(--cyan); animation-delay: 1.5s; }
.hud-corner.br { bottom: 60px; right: 20px; border-bottom: 2px solid var(--cyan); border-right: 2px solid var(--cyan); animation-delay: 2.25s; }

.master-badge-wrap {
    position: relative; display: inline-flex; align-items: center; justify-content: center;
    margin-bottom: 1.6rem;
    animation: slide-up 0.5s ease-out backwards, float-y 4s ease-in-out infinite 0.8s;
    padding: 0 8px;
}
.master-badge {
    position: relative; display: inline-flex; align-items: center; gap: 14px;
    padding: 0.85rem 2.2rem; border-radius: 6px;
    font-family: 'Share Tech Mono', monospace;
    font-size: clamp(0.85rem, 2.4vw, 1.35rem); letter-spacing: 5px; text-transform: uppercase;
    color: var(--ink-bright); background: var(--bg-surface); border: 1px solid transparent;
    z-index: 1; animation: mega-glow 3s ease-in-out infinite; overflow: hidden;
}
.master-badge-wrap::before {
    content: ""; position: absolute; inset: -2px; border-radius: 8px;
    background: linear-gradient(90deg, var(--cyan), var(--purple), var(--green), var(--cyan));
    background-size: 200% 100%; animation: border-rotate 3s linear infinite; z-index: 0;
}
.master-badge-wrap::after {
    content: ""; position: absolute; inset: 0; border-radius: 6px;
    background: var(--bg-surface); z-index: 0; margin: 1px;
}
.master-badge-icon {
    color: var(--cyan); font-size: 1.1em;
    animation: rotate-slow 5s linear infinite; display: inline-block;
    text-shadow: 0 0 14px rgba(0,229,255,0.7); position: relative; z-index: 1;
}
.master-badge-text-1, .master-badge-text-2, .master-badge-text-3, .master-badge-amp {
    position: relative; z-index: 1;
}
.master-badge-text-1 { color: var(--ink-bright); }
.master-badge-text-2 { color: var(--cyan); text-shadow: 0 0 16px rgba(0,229,255,0.6); }
.master-badge-text-3 { color: var(--purple); text-shadow: 0 0 16px rgba(124,58,255,0.5); }
.master-badge-amp { color: var(--green); text-shadow: 0 0 12px rgba(0,255,136,0.6); animation: dot-blink 2s ease-in-out infinite; }
.master-badge::before {
    content: ""; position: absolute; top: 0; left: 0; right: 0; bottom: 0;
    background: linear-gradient(110deg, transparent 30%, rgba(0,229,255,0.18) 50%, transparent 70%);
    background-size: 200% 100%; animation: badge-shimmer 3s linear infinite;
    border-radius: 6px; pointer-events: none; z-index: 0;
}

.system-badge {
    display: inline-flex; align-items: center; gap: 8px;
    background: transparent; border: 1px solid var(--cyan-border);
    border-radius: 2px; padding: 0.4rem 1.2rem;
    font-family: 'Share Tech Mono', monospace; font-size: 0.65rem;
    color: var(--cyan); letter-spacing: 3px; margin-bottom: 1.4rem; text-transform: uppercase;
    animation: slide-up 0.5s ease-out, pulse-cyan 4s ease-in-out infinite 1s;
    cursor: default;
}
.system-badge .badge-dot {
    width: 5px; height: 5px; border-radius: 50%; background: var(--cyan);
    animation: dot-blink 1.5s ease-in-out infinite;
}

.hero-title {
    font-family: 'Share Tech Mono', monospace;
    font-size: clamp(2.4rem, 7vw, 5rem); line-height: 1.05;
    color: var(--ink-bright); margin-bottom: 0.1rem; letter-spacing: 2px;
    animation: slide-up 0.55s ease-out 0.1s backwards; cursor: default;
}
.hero-title-accent {
    font-family: 'Share Tech Mono', monospace;
    font-size: clamp(2.2rem, 6.5vw, 4.6rem); color: var(--cyan); display: block;
    margin-bottom: 1.6rem; letter-spacing: 1px;
    text-shadow: 0 0 30px rgba(0,229,255,0.45), 0 0 60px rgba(0,229,255,0.15);
    animation: slide-up 0.6s ease-out 0.2s backwards, glow-pulse 3s ease-in-out infinite 1s;
    cursor: default;
}
.cursor-blink {
    display: inline-block; width: 3px; height: 0.85em; background: var(--cyan);
    margin-left: 4px; vertical-align: middle; animation: blink-cursor 1s step-end infinite;
}

.hero-desc {
    font-family: 'JetBrains Mono', monospace;
    font-size: clamp(0.78rem, 1.5vw, 0.9rem); color: var(--ink-muted);
    max-width: 600px; margin: 0 auto 2rem auto; line-height: 1.8; font-weight: 400;
    animation: slide-up 0.65s ease-out 0.25s backwards;
}
.hero-desc strong { color: var(--ink); }

.creator-card {
    background: var(--bg-surface); border: 1px solid var(--cyan-border);
    border-left: 3px solid var(--cyan); border-radius: 4px;
    padding: 1rem 1.4rem; max-width: 500px; margin: 0 auto 2rem auto;
    display: flex; align-items: center; gap: 1rem; text-align: left;
    animation: slide-up 0.7s ease-out 0.3s backwards, creator-float 6s ease-in-out infinite 1.5s;
    position: relative; overflow: hidden; cursor: default;
}
.creator-card::before {
    content: ""; position: absolute; inset: -2px; border-radius: 6px;
    background: linear-gradient(120deg, transparent, var(--cyan), transparent 40%);
    background-size: 250% 250%; opacity: 0.4;
    animation: border-rotate 4s linear infinite; z-index: 0;
}
.creator-card::after {
    content: ""; position: absolute; inset: 1px; border-radius: 4px;
    background: var(--bg-surface); z-index: 0;
}
.creator-card > * { position: relative; z-index: 1; }
.creator-avatar-wrap {
    position: relative; width: 48px; height: 48px; min-width: 48px;
    display: flex; align-items: center; justify-content: center;
}
.creator-avatar-wrap::before, .creator-avatar-wrap::after {
    content: ""; position: absolute; inset: -6px; border-radius: 8px; border: 1px dashed var(--cyan-border);
}
.creator-avatar-wrap::before { animation: ring-spin 6s linear infinite; }
.creator-avatar-wrap::after { inset: -11px; border: 1px dotted var(--purple-dim); animation: ring-spin-reverse 10s linear infinite; }
.creator-avatar {
    width: 48px; height: 48px; min-width: 48px; border-radius: 4px;
    background: linear-gradient(135deg, var(--purple), #3B5FE0);
    display: flex; align-items: center; justify-content: center;
    font-family: 'Share Tech Mono', monospace; font-size: 1rem; color: #fff;
    border: 1px solid rgba(124,58,255,0.5);
    animation: avatar-pulse 2.6s ease-in-out infinite; position: relative; z-index: 1;
}
.creator-name {
    font-family: 'Share Tech Mono', monospace; font-size: 0.9rem;
    color: var(--ink-bright); letter-spacing: 1px; margin-bottom: 2px;
    animation: glow-pulse 4s ease-in-out infinite;
}
.creator-role { font-family: 'JetBrains Mono', monospace; font-size: 0.65rem; color: var(--ink-muted); letter-spacing: 1px; }
.creator-tags { display: flex; gap: 0.4rem; margin-top: 0.4rem; flex-wrap: wrap; }
.ctag {
    font-family: 'JetBrains Mono', monospace; font-size: 0.58rem; padding: 2px 8px;
    border: 1px solid; border-radius: 2px; letter-spacing: 1px; text-transform: uppercase;
    font-weight: 600; animation: tag-glow-cycle 3s ease-in-out infinite;
}
.ctag-cyan  { border-color: var(--cyan-dim); color: var(--cyan-dim); animation-delay: 0s; }
.ctag-green { border-color: rgba(0,255,136,0.4); color: var(--green); animation-delay: 1s; }
.ctag-purp  { border-color: rgba(124,58,255,0.4); color: var(--purple); animation-delay: 2s; }

.flow-strip {
    display: flex; align-items: stretch; gap: 0; max-width: 860px;
    margin: 0 auto 2rem auto; width: 100%;
    animation: slide-up 0.7s ease-out 0.35s backwards;
}
.flow-card {
    background: var(--bg-surface); border: 1px solid var(--cyan-border);
    padding: 1.4rem 1.2rem; flex: 1; position: relative; overflow: hidden;
    transition: all 0.3s ease; cursor: default;
}
.flow-card:first-child {
    border-top: 2px solid var(--cyan); border-radius: 6px 0 0 6px;
    animation: flow-card-lift 4s ease-in-out infinite 0s;
}
.flow-card:nth-child(3) {
    border-top: 2px solid var(--green);
    animation: flow-card-lift 4s ease-in-out infinite 1.3s;
}
.flow-card:last-child {
    border-top: 2px solid var(--purple); border-radius: 0 6px 6px 0;
    animation: flow-card-lift 4s ease-in-out infinite 2.6s;
}
.flow-card::before {
    content: ""; position: absolute; top: 0; left: -60%; width: 40%; height: 100%;
    background: linear-gradient(120deg, transparent, rgba(0,229,255,0.10), transparent);
    transform: skewX(-20deg); animation: badge-shimmer 4s linear infinite;
    pointer-events: none;
}
.flow-num { font-family: 'Share Tech Mono', monospace; font-size: 0.6rem; color: var(--ink-muted); letter-spacing: 2px; margin-bottom: 0.6rem; display: block; }
.flow-icon {
    font-size: 1.8rem; margin-bottom: 0.6rem;
    display: inline-block; animation: rotate-slow 8s linear infinite;
    transition: transform 0.35s ease;
}
.flow-title {
    font-family: 'Share Tech Mono', monospace; font-size: 0.85rem; color: var(--ink-bright);
    margin-bottom: 0.4rem; letter-spacing: 1px; text-transform: uppercase;
    animation: glow-pulse 5s ease-in-out infinite;
}
.flow-desc { font-family: 'JetBrains Mono', monospace; font-size: 0.72rem; color: var(--ink-muted); line-height: 1.55; }
.flow-sep { width: 1px; background: var(--cyan-border); align-self: stretch; }

.stat-strip {
    display: flex; gap: 1px; justify-content: center; max-width: 680px;
    margin: 0 auto 2rem auto; width: 100%;
    background: var(--cyan-border); border: 1px solid var(--cyan-border);
    border-radius: 4px; overflow: hidden;
    animation: slide-up 0.7s ease-out 0.4s backwards;
}
.stat-pill {
    background: var(--bg-surface); padding: 1rem 1.4rem; text-align: center; flex: 1;
    cursor: default;
}
.stat-pill:nth-child(1) { animation: pulse-cyan 3s ease-in-out infinite 0s; }
.stat-pill:nth-child(2) { animation: pulse-cyan 3s ease-in-out infinite 0.75s; }
.stat-pill:nth-child(3) { animation: pulse-cyan 3s ease-in-out infinite 1.5s; }
.stat-pill:nth-child(4) { animation: pulse-cyan 3s ease-in-out infinite 2.25s; }
.stat-num {
    font-family: 'Share Tech Mono', monospace; font-size: 1.3rem; color: var(--cyan);
    display: inline-block;
    text-shadow: 0 0 14px rgba(0,229,255,0.4); animation: glow-pulse 3s ease-in-out infinite;
}
.stat-label {
    font-family: 'JetBrains Mono', monospace; font-size: 0.55rem; color: var(--ink-muted);
    letter-spacing: 2px; text-transform: uppercase; margin-top: 2px; display: block;
}

.hero-stack {
    display: flex; align-items: center; gap: 0.5rem; flex-wrap: wrap;
    justify-content: center; animation: slide-up 0.7s ease-out 0.45s backwards;
}
.stack-pill {
    background: transparent; border: 1px solid var(--ink-faint);
    border-radius: 2px; padding: 0.3rem 0.8rem;
    font-family: 'JetBrains Mono', monospace; font-size: 0.65rem; color: var(--ink-muted);
    letter-spacing: 1px; text-transform: uppercase;
    font-weight: 500; position: relative; overflow: hidden;
}
.stack-pill:nth-child(1) { animation: pill-glow 3s ease-in-out infinite 0s; }
.stack-pill:nth-child(2) { animation: pill-glow 3s ease-in-out infinite 0.6s; }
.stack-pill:nth-child(3) { animation: pill-glow 3s ease-in-out infinite 1.2s; }
.stack-pill:nth-child(4) { animation: pill-glow 3s ease-in-out infinite 1.8s; }
.stack-pill:nth-child(5) { animation: pill-glow 3s ease-in-out infinite 2.4s; }

.main .stButton > button[kind="primary"],
[data-testid="stMain"] .stButton > button[kind="primary"] {
    background: transparent !important; border: 1px solid var(--cyan) !important;
    color: var(--cyan) !important; font-family: 'Share Tech Mono', monospace !important;
    font-size: 0.82rem !important; font-weight: 400 !important; letter-spacing: 3px !important;
    border-radius: 4px !important; padding: 14px 28px !important;
    box-shadow: 0 0 16px rgba(0,229,255,0.1) !important; transition: all 0.15s ease !important;
    text-transform: uppercase !important;
    animation: pulse-cyan 3s ease-in-out infinite !important;
}
.main .stButton > button[kind="primary"]:hover,
[data-testid="stMain"] .stButton > button[kind="primary"]:hover {
    background: rgba(0,229,255,0.08) !important;
    box-shadow: 0 0 30px rgba(0,229,255,0.2), inset 0 0 20px rgba(0,229,255,0.05) !important;
    text-shadow: 0 0 8px rgba(0,229,255,0.8) !important;
}

.main .stButton > button,
[data-testid="stMain"] .stButton > button {
    background: transparent !important; border: 1px solid var(--ink-faint) !important;
    color: var(--ink-muted) !important; font-family: 'JetBrains Mono', monospace !important;
    font-size: 0.75rem !important; font-weight: 500 !important; border-radius: 4px !important;
    padding: 10px 20px !important; box-shadow: none !important; transition: all 0.15s ease !important;
    text-transform: uppercase !important; letter-spacing: 1.5px !important;
}
.main .stButton > button:hover,
[data-testid="stMain"] .stButton > button:hover {
    background: var(--cyan-glow) !important; border-color: var(--cyan-border) !important; color: var(--ink) !important;
}
.main .stButton > button:focus,
[data-testid="stMain"] .stButton > button:focus,
.main .stButton > button[kind="primary"]:focus,
[data-testid="stMain"] .stButton > button[kind="primary"]:focus { box-shadow: none !important; outline: none !important; }

.stDownloadButton > button {
    background: transparent !important; border: 1px solid var(--green) !important;
    color: var(--green) !important; font-family: 'Share Tech Mono', monospace !important;
    font-size: 0.8rem !important; font-weight: 400 !important; letter-spacing: 2px !important;
    border-radius: 4px !important; padding: 13px 24px !important;
    box-shadow: 0 0 12px rgba(0,255,136,0.1) !important; transition: all 0.15s ease !important;
    text-transform: uppercase !important;
}
.stDownloadButton > button:hover { background: var(--green-dim) !important; box-shadow: 0 0 24px rgba(0,255,136,0.2) !important; }

.stLinkButton > a {
    background: transparent !important; border: 1px solid var(--ink-faint) !important;
    color: var(--ink-muted) !important; font-family: 'JetBrains Mono', monospace !important;
    font-size: 0.72rem !important; font-weight: 500 !important; border-radius: 4px !important;
    padding: 10px 18px !important; box-shadow: none !important; transition: all 0.15s ease !important;
    letter-spacing: 1px !important; text-transform: uppercase !important;
}
.stLinkButton > a:hover { border-color: var(--cyan-border) !important; color: var(--cyan-dim) !important; background: var(--cyan-glow) !important; }

.stTextInput > div > div > input,
.stTextArea > div > div > textarea {
    background: var(--bg-raised) !important; border: 1px solid var(--ink-faint) !important;
    border-radius: 4px !important; color: var(--ink) !important;
    font-family: 'JetBrains Mono', monospace !important; font-size: 0.85rem !important;
    transition: all 0.15s ease !important; caret-color: var(--cyan) !important;
}
.stTextInput > div > div > input:focus,
.stTextArea > div > div > textarea:focus {
    border-color: var(--cyan-dim) !important;
    box-shadow: 0 0 14px rgba(0,229,255,0.08), inset 0 0 8px rgba(0,229,255,0.03) !important;
}
.stTextInput > div > div > input::placeholder,
.stTextArea > div > div > textarea::placeholder { color: var(--ink-faint) !important; opacity: 1 !important; }
.stTextInput label, .stTextArea label, .stSelectbox label, .stNumberInput label {
    font-family: 'Share Tech Mono', monospace !important; color: var(--ink-muted) !important;
    font-size: 0.7rem !important; font-weight: 400 !important;
    letter-spacing: 2px !important; text-transform: uppercase !important;
}
[data-baseweb="select"] > div {
    background: var(--bg-raised) !important; border: 1px solid var(--ink-faint) !important;
    border-radius: 4px !important; color: var(--ink) !important;
    font-family: 'JetBrains Mono', monospace !important; font-size: 0.82rem !important;
    transition: border-color 0.15s ease !important;
}
[data-baseweb="select"] > div:hover { border-color: var(--cyan-dim) !important; }
[data-baseweb="popover"] { background: var(--bg-raised) !important; border: 1px solid var(--cyan-border) !important; }
[data-baseweb="menu"] { background: var(--bg-raised) !important; }
[data-baseweb="option"] { background: var(--bg-raised) !important; color: var(--ink) !important; font-family: 'JetBrains Mono', monospace !important; font-size: 0.82rem !important; }
[data-baseweb="option"]:hover { background: var(--bg-hover) !important; color: var(--cyan) !important; }

[data-testid="stRadio"] > div { gap: 0.5rem !important; }
[data-testid="stRadio"] label {
    background: var(--bg-surface) !important; border: 1px solid var(--ink-faint) !important;
    border-radius: 4px !important; padding: 0.45rem 1.1rem !important; transition: all 0.15s ease !important; cursor: pointer !important;
}
[data-testid="stRadio"] label:hover { border-color: var(--cyan-dim) !important; background: var(--cyan-glow) !important; }
[data-testid="stRadio"] label p, [data-testid="stRadio"] div {
    color: var(--ink) !important; font-family: 'JetBrains Mono', monospace !important;
    font-weight: 400 !important; font-size: 0.8rem !important; letter-spacing: 0.5px !important;
}

.stProgress > div > div > div > div {
    background: linear-gradient(90deg, var(--cyan-dim), var(--cyan)) !important;
    border-radius: 0 !important; box-shadow: 0 0 10px rgba(0,229,255,0.4) !important;
}
.stProgress > div > div { background: var(--bg-raised) !important; border-radius: 0 !important; border: 1px solid var(--ink-faint) !important; height: 6px !important; }
.stProgress div[data-testid="stProgressBar"] span { color: var(--ink-muted) !important; font-family: 'JetBrains Mono', monospace !important; font-size: 0.7rem !important; }

[data-testid="metric-container"] {
    background: var(--bg-surface) !important; border: 1px solid var(--cyan-border) !important;
    border-radius: 4px !important; padding: 14px !important; text-align: center !important;
    transition: all 0.18s ease !important; box-shadow: 0 0 14px rgba(0,229,255,0.04) !important;
}
[data-testid="metric-container"]:hover { border-color: var(--cyan-strong) !important; box-shadow: 0 0 20px var(--cyan-glow) !important; }
[data-testid="stMetricLabel"] { font-family: 'Share Tech Mono', monospace !important; color: var(--ink-muted) !important; font-size: 0.65rem !important; letter-spacing: 2px !important; text-transform: uppercase !important; }
[data-testid="stMetricValue"] { font-family: 'Share Tech Mono', monospace !important; color: var(--cyan) !important; font-size: 1.8rem !important; text-shadow: 0 0 16px rgba(0,229,255,0.4) !important; }

[data-testid="stExpander"] { border: 1px solid var(--cyan-border) !important; border-radius: 4px !important; background: var(--bg-surface) !important; box-shadow: none !important; }
[data-testid="stExpander"] summary p { color: var(--ink) !important; font-family: 'JetBrains Mono', monospace !important; font-weight: 500 !important; font-size: 0.8rem !important; letter-spacing: 1px !important; }
[data-testid="stExpander"]:hover { border-color: var(--cyan-strong) !important; }

[data-testid="stSpinner"] p, [data-testid="stSpinner"] div { color: var(--cyan-dim) !important; font-family: 'JetBrains Mono', monospace !important; font-size: 0.8rem !important; letter-spacing: 1px !important; }

.stAlert { background: var(--bg-surface) !important; border-radius: 4px !important; border: 1px solid var(--cyan-border) !important; box-shadow: none !important; font-family: 'JetBrains Mono', monospace !important; font-size: 0.8rem !important; }
.stAlert p, .stAlert div, .stAlert span { color: var(--ink) !important; font-family: 'JetBrains Mono', monospace !important; }
.stCodeBlock { background: var(--bg-raised) !important; border: 1px solid var(--cyan-border) !important; border-radius: 4px !important; }
.stMarkdown p, .stMarkdown li { font-family: 'JetBrains Mono', monospace !important; color: var(--ink-muted) !important; font-size: 0.82rem !important; }
.stMarkdown strong { color: var(--cyan) !important; }
hr { border: none !important; border-top: 1px solid var(--ink-faint) !important; margin: 1.2rem 0 !important; }

.q-card {
    background: var(--bg-surface); border: 1px solid var(--cyan-border);
    border-left: 3px solid var(--purple); border-radius: 4px;
    padding: 1rem 1.3rem; margin-bottom: 0.8rem; transition: all 0.18s ease;
    animation: slide-up 0.3s ease-out backwards; position: relative;
}
.q-card:hover { border-color: var(--cyan-strong); border-left-color: var(--cyan); box-shadow: 0 0 18px var(--cyan-glow); }
.q-num-badge {
    display: inline-flex; align-items: center; justify-content: center;
    width: 22px; height: 22px; border-radius: 2px;
    background: var(--purple-dim); border: 1px solid rgba(124,58,255,0.4);
    color: var(--purple); font-family: 'Share Tech Mono', monospace; font-size: 0.68rem;
    margin-right: 8px; flex-shrink: 0;
}
.q-title { font-family: 'JetBrains Mono', monospace; color: var(--ink); font-size: 0.88rem; font-weight: 500; margin-bottom: 0.7rem; display: flex; align-items: center; }
.q-opt { padding: 5px 10px; margin: 3px 0; font-family: 'JetBrains Mono', monospace; font-size: 0.8rem; color: var(--ink-muted); border-radius: 2px; background: transparent; border: 1px solid var(--ink-faint); transition: all 0.12s ease; }
.q-opt:hover { background: var(--cyan-glow); border-color: var(--cyan-border); color: var(--ink); }
.q-correct { padding: 6px 10px; margin: 3px 0; font-family: 'JetBrains Mono', monospace; font-size: 0.8rem; color: var(--green); font-weight: 500; border-radius: 2px; background: var(--green-dim); border: 1px solid rgba(0,255,136,0.3); }
.q-exp { font-size: 0.75rem; color: var(--ink-muted); margin-top: 0.7rem; font-family: 'JetBrains Mono', monospace; background: var(--bg-raised); border-left: 2px solid var(--amber); padding: 0.4rem 0.8rem; border-radius: 0 3px 3px 0; }
.ak-badge { display: inline-block; background: transparent; border: 1px solid rgba(0,229,255,0.3); border-radius: 2px; padding: 1px 8px; font-family: 'Share Tech Mono', monospace; font-size: 0.58rem; color: var(--cyan-dim); margin-left: 8px; letter-spacing: 1px; text-transform: uppercase; }

.result-row {
    display: flex; align-items: center; gap: 1rem; background: var(--bg-surface);
    border: 1px solid var(--ink-faint); border-radius: 3px; padding: 0.6rem 1rem;
    margin-bottom: 0.4rem; transition: all 0.15s ease; animation: slide-right 0.35s ease-out backwards;
    font-family: 'JetBrains Mono', monospace;
}
.result-row:hover { border-color: var(--cyan-border); background: var(--bg-hover); }
.result-row.pass { border-left: 3px solid var(--green); }
.result-row.fail { border-left: 3px solid var(--red); }
.grade-pill {
    display: inline-flex; align-items: center; justify-content: center;
    min-width: 40px; height: 28px; border-radius: 2px;
    font-family: 'Share Tech Mono', monospace; font-size: 0.8rem; padding: 0 6px; border: 1px solid; letter-spacing: 1px;
}
.grade-pill.pass { background: var(--green-dim); border-color: rgba(0,255,136,0.4); color: var(--green); }
.grade-pill.fail { background: var(--red-dim); border-color: rgba(255,68,68,0.4); color: var(--red); }

.celebrate-banner {
    background: var(--bg-surface); border: 1px solid var(--cyan-border);
    border-top: 2px solid var(--cyan); border-radius: 4px; padding: 1.2rem 1.4rem;
    text-align: center; margin-bottom: 1rem; position: relative; overflow: hidden;
}
.celebrate-banner::after {
    content: ""; position: absolute; top: -50%; left: -50%;
    width: 20px; height: 200%;
    background: linear-gradient(90deg, transparent, rgba(0,229,255,0.08), transparent);
    animation: scan-line 4s linear infinite;
}
.celebrate-title { font-family: 'Share Tech Mono', monospace; font-size: 1.25rem; color: var(--cyan); margin-bottom: 0.3rem; letter-spacing: 2px; text-shadow: 0 0 20px rgba(0,229,255,0.5); animation: glow-pulse 2.5s ease-in-out infinite; display: inline-block; }
.celebrate-sub { font-family: 'JetBrains Mono', monospace; color: var(--ink-muted); font-size: 0.75rem; letter-spacing: 1px; }

.empty-state { text-align: center; padding: 2.8rem 2rem; font-family: 'JetBrains Mono', monospace; color: var(--ink-muted); font-size: 0.82rem; letter-spacing: 1px; background: var(--bg-surface); border-radius: 4px; border: 1px dashed var(--cyan-border); }
.empty-state .empty-emoji { font-size: 2rem; display: block; margin-bottom: 0.7rem; opacity: 0.6; }

.dropzone { text-align: center; padding: 2.4rem 1.5rem; border: 1px dashed var(--cyan-border); border-radius: 4px; color: var(--ink-muted); font-family: 'JetBrains Mono', monospace; font-size: 0.78rem; letter-spacing: 1px; background: var(--bg-raised); transition: all 0.2s ease; }
.dropzone:hover { border-color: var(--cyan-strong); background: var(--bg-hover); }
.dropzone .dz-icon { font-size: 2rem; display: block; margin-bottom: 0.6rem; opacity: 0.5; }

@media (max-width: 760px) {
    .journey-bar { flex-direction: column; align-items: flex-start; gap: 0.8rem; }
    .journey-connector { display: none; }
    .flow-strip { flex-direction: column; }
    .flow-sep { display: none; }
    .flow-card:first-child { border-radius: 6px 6px 0 0; }
    .flow-card:last-child  { border-radius: 0 0 6px 6px; }
}
@media (max-width: 640px) {
    .master-badge { letter-spacing: 2px; padding: 0.7rem 1.2rem; gap: 8px; flex-wrap: wrap; justify-content: center; }
}
</style>
""", unsafe_allow_html=True)

SCOPES = [
    "https://www.googleapis.com/auth/forms.body",
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/drive",
]

try:
    api_key = st.secrets["GROQ_API_KEY"]
except Exception:
    api_key = os.getenv("GROQ_API_KEY")

if not api_key:
    st.error("GROQ_API_KEY not found — add it to your secrets.")
    st.stop()
client = Groq(api_key=api_key)


# ══════════════════════════════════════════════════════
#  IMPROVEMENT #1 — Smart content truncation
# ══════════════════════════════════════════════════════
def smart_truncate(text, max_chars=12000):
    """Cut at the last complete sentence within max_chars to avoid mid-sentence context."""
    if len(text) <= max_chars:
        return text
    truncated = text[:max_chars]
    last_stop = max(truncated.rfind('.'), truncated.rfind('?'), truncated.rfind('!'))
    if last_stop > max_chars * 0.7:
        return truncated[:last_stop + 1] + "\n[...content truncated for length...]"
    return truncated + "\n[...content truncated for length...]"


# ══════════════════════════════════════════════════════
#  IMPROVEMENT #2 — Robust JSON extractor + validator
# ══════════════════════════════════════════════════════
def parse_and_validate_questions(raw_text, expected_count):
    """
    Strip markdown fences, extract the JSON array wherever it appears,
    parse it, and validate every question has the required fields.
    Raises ValueError with a clear message on any failure.
    """
    # Remove any markdown code fences the model may have added
    cleaned = re.sub(r'```(?:json)?\s*', '', raw_text).strip()
    cleaned = re.sub(r'```', '', cleaned).strip()

    # Extract the outermost JSON array even if the model added preamble/postamble
    match = re.search(r'\[.*\]', cleaned, re.DOTALL)
    if not match:
        raise ValueError(
            "Model did not return a JSON array. "
            "The response may have been cut off or the model added unexpected text."
        )

    questions = json.loads(match.group(0))

    if not isinstance(questions, list) or len(questions) == 0:
        raise ValueError("Parsed result is empty or not a list.")

    required_keys = {"question", "options", "answer", "explanation"}
    for i, q in enumerate(questions):
        missing = required_keys - set(q.keys())
        if missing:
            raise ValueError(f"Question {i+1} is missing fields: {missing}")
        if not isinstance(q["options"], list) or len(q["options"]) != 4:
            raise ValueError(f"Question {i+1} must have exactly 4 options (got {len(q.get('options', []))}).")
        if q["answer"] not in q["options"]:
            raise ValueError(
                f"Question {i+1}: answer '{q['answer']}' not found in its options list. "
                "The model may have formatted the answer differently from the options."
            )

    return questions


# ══════════════════════════════════════════════════════
#  IMPROVEMENT #3 — High-quality prompt builder
# ══════════════════════════════════════════════════════
def build_prompt(subject, difficulty, num_q, content):
    diff_map = {
        "Easy":   (
            "Recall and basic comprehension. Test definitions, simple facts, and "
            "direct application of a single concept. Distractors should be clearly "
            "wrong to anyone who read the material."
        ),
        "Medium": (
            "Application and analysis. Students must apply or compare 2+ concepts. "
            "Use plausible distractors that require genuine understanding to reject."
        ),
        "Hard":   (
            "Synthesis and evaluation. Require multi-step reasoning, edge cases, or "
            "nuanced distinctions. All 4 options must be plausible — only a student "
            "who deeply understands the material can identify the correct one."
        ),
        "Mixed":  (
            "Distribute difficulty evenly: roughly one third Easy (recall), "
            "one third Medium (application), one third Hard (synthesis). "
            "Vary the order randomly so the quiz doesn't cluster by difficulty."
        ),
    }
    diff_instruction = diff_map.get(difficulty, diff_map["Medium"])

    return f"""You are a senior university professor and psychometrician writing a high-stakes academic quiz.

SUBJECT: {subject}
DIFFICULTY: {difficulty}
DIFFICULTY GUIDE: {diff_instruction}
TOTAL QUESTIONS REQUIRED: exactly {num_q}

SOURCE MATERIAL:
\"\"\"
{content}
\"\"\"

STRICT AUTHORING RULES — violating any rule makes the quiz unusable:
1. Every question must be answerable from the SOURCE MATERIAL above. No outside knowledge.
2. Each question must test a DISTINCT concept — absolutely no repetition across questions.
3. Distractors (wrong options) must be plausible but unambiguously wrong to a prepared student.
4. The correct answer must be provably correct with no reasonable alternative interpretation.
5. Explanations must state WHY the correct answer is right AND why the most tempting distractor is wrong.
6. Vary question formats across the set: definitions, application, comparison, cause-and-effect,
   "which of the following is NOT true", scenario-based, fill-in-the-blank style, etc.
7. Never start two questions with the same opening word or phrase.
8. Never use "All of the above" or "None of the above" as an option.
9. Keep option lengths roughly equal so length alone cannot hint at the answer.
10. Number each option exactly as shown: "A) ...", "B) ...", "C) ...", "D) ..."

OUTPUT FORMAT — return ONLY a raw JSON array. No markdown fences, no preamble, no trailing text.
Your response must start with [ and end with ].

EXACT SCHEMA per question:
{{
  "question": "Full question text ending with a question mark?",
  "options": ["A) First option", "B) Second option", "C) Third option", "D) Fourth option"],
  "answer": "A) First option",
  "explanation": "A is correct because [specific reason from the material]. B is tempting but wrong because [specific reason].",
  "difficulty_tag": "Easy",
  "concept": "short concept name this question tests"
}}

Generate exactly {num_q} question objects now. Start your response with [ and end with ]."""


# ══════════════════════════════════════════════════════
#  IMPROVEMENT #4 — Retry wrapper with exponential back-off
# ══════════════════════════════════════════════════════
def generate_questions_with_retry(subject, difficulty, num_q, content, max_retries=3):
    """
    Call the Groq API with the best available model and retry up to max_retries
    times on parse/validation errors.

    Model: llama-3.3-70b-versatile  (current supported successor to 3.1 70b)
    Temperature: 0.4  (lower = more factually consistent)
    max_tokens: 8192  (large enough for 50 questions with explanations)
    """
    prompt = build_prompt(subject, difficulty, num_q, content)

    system_msg = (
        "You are a senior university professor and expert psychometrician. "
        "Your sole task is to produce high-quality multiple-choice questions "
        "in strict JSON format. Rules you never break: "
        "(1) Output ONLY a raw JSON array — no markdown, no explanation, no preamble. "
        "(2) Start your response with [ and end with ]. "
        "(3) Every answer must be provably correct from the given source material. "
        "(4) All distractors must be plausible but definitively wrong. "
        "(5) Never repeat a concept. "
        "(6) Never use 'All of the above' or 'None of the above'."
    )

    last_error = None
    for attempt in range(1, max_retries + 1):
        try:
            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[
                    {"role": "system", "content": system_msg},
                    {"role": "user",   "content": prompt},
                ],
                max_tokens=8192,
                temperature=0.4,
            )
            raw = resp.choices[0].message.content.strip()
            questions = parse_and_validate_questions(raw, num_q)
            return questions

        except json.JSONDecodeError as e:
            last_error = f"JSON parse error on attempt {attempt}: {e}"
        except ValueError as e:
            last_error = f"Validation error on attempt {attempt}: {e}"
        except Exception as e:
            last_error = f"API error on attempt {attempt}: {e}"

        if attempt < max_retries:
            time.sleep(attempt * 1.5)

    raise RuntimeError(
        f"Generation failed after {max_retries} attempts. Last error: {last_error}"
    )

def render_quest_tracker(active_key):
    order = ["step1", "step2", "step3"]
    icons = ["01", "02", "03"]
    labels = ["CREATE", "PUBLISH", "RESULTS"]
    try:
        active_idx = order.index(active_key)
    except ValueError:
        active_idx = -1
    html = '<div class="quest-tracker">'
    for i, (key, icon, label) in enumerate(zip(order, icons, labels)):
        if i == active_idx:
            circle_cls, label_cls, content = "active", "active", icon
        elif i < active_idx:
            circle_cls, label_cls, content = "done", "", "✓"
        else:
            circle_cls, label_cls, content = "", "", icon
        html += f'<div class="quest-node"><div class="quest-circle {circle_cls}">{content}</div><div class="quest-label {label_cls}">{label}</div></div>'
        if i < len(order) - 1:
            line_cls = "done" if i < active_idx else ""
            html += f'<div class="quest-line {line_cls}"></div>'
    html += '</div>'
    st.markdown(html, unsafe_allow_html=True)


def render_journey_bar(active_key):
    order = ["step1", "step2", "step3"]
    icons = ["01", "02", "03"]
    names = ["CREATE_QUIZ", "PUBLISH_FORM", "CHECK_RESULTS"]
    subs = ["AI generates questions", "Deploy to Google Forms", "Auto-grade & export"]
    try:
        active_idx = order.index(active_key)
    except ValueError:
        active_idx = -1
    html = '<div class="journey-bar">'
    for i, (key, icon, name, sub) in enumerate(zip(order, icons, names, subs)):
        if i == active_idx:
            circle_cls, content = "active", icon
        elif i < active_idx:
            circle_cls, content = "done", "✓"
        else:
            circle_cls, content = "", icon
        html += f'''<div class="journey-step">
            <div class="journey-circle {circle_cls}">{content}</div>
            <div class="journey-text">
                <span class="journey-name">{name}</span>
                <span class="journey-sub">{sub}</span>
            </div>
        </div>'''
        if i < len(order) - 1:
            conn_cls = "done" if i < active_idx else ""
            html += f'<div class="journey-connector {conn_cls}"></div>'
    html += '</div>'
    st.markdown(html, unsafe_allow_html=True)


def get_grade(pct):
    if pct >= 90: return "A+"
    if pct >= 80: return "A"
    if pct >= 70: return "B"
    if pct >= 60: return "C"
    if pct >= 50: return "D"
    return "F"

def get_grade_emoji(pct):
    if pct >= 90: return "◆"
    if pct >= 80: return "▲"
    if pct >= 70: return "●"
    if pct >= 60: return "○"
    if pct >= 50: return "△"
    return "▽"

def make_excel(results, subject, total):
    wb = Workbook(); ws = wb.active; ws.title = "Quiz Results"
    h_fill = PatternFill("solid", fgColor="060D0F")
    h_font = Font(color="00E5FF", bold=True, size=11)
    center = Alignment(horizontal="center", vertical="center")
    thin = Side(style="thin", color="0C1A1E")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    headers = ["#","Student Name","Roll Number","Marks Obtained",f"Total ({total})","Percentage","Grade","Status"]
    ws.append(headers)
    for c in range(1, len(headers)+1):
        cell = ws.cell(1, c)
        cell.fill = h_fill; cell.font = h_font; cell.alignment = center; cell.border = border
    ws.row_dimensions[1].height = 22
    for i, w in enumerate([5,22,18,16,14,12,10,10], 1):
        ws.column_dimensions[ws.cell(1,i).column_letter].width = w
    pass_fill = PatternFill("solid", fgColor="E6FFF4")
    fail_fill = PatternFill("solid", fgColor="FFE6E6")
    for idx, s in enumerate(results, 1):
        pct = round((s["marks"]/total)*100, 1)
        grade = get_grade(pct); status = "Pass" if pct >= 50 else "Fail"
        ws.append([idx, s["name"], s["roll"], s["marks"], total, f"{pct}%", grade, status])
        fill = pass_fill if status == "Pass" else fail_fill
        for c in range(1, len(headers)+1):
            cell = ws.cell(idx+1, c); cell.alignment = center; cell.border = border; cell.fill = fill
    ws.append([])
    passed = sum(1 for s in results if (s["marks"]/total)*100 >= 50)
    avg = round(sum(s["marks"] for s in results)/len(results), 1) if results else 0
    sr = ws.max_row + 1
    bold = Font(bold=True, size=11)
    ws.cell(sr,1,"Summary").font = bold
    ws.cell(sr,2,f"Total Students: {len(results)}").font = bold
    ws.cell(sr,3,f"Passed: {passed}").font = Font(bold=True, color="059669")
    ws.cell(sr,4,f"Failed: {len(results)-passed}").font = Font(bold=True, color="DC2626")
    ws.cell(sr,5,f"Class Avg: {avg}/{total}").font = bold
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf.read()

def extract_text_from_pdf(uploaded_file):
    try:
        pdf_bytes = uploaded_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        text = ""
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text += f"\n--- Page {page_num+1} ---\n" + page.get_text()
        doc.close()
        return text.strip()
    except Exception as e:
        return f"ERROR: {str(e)}"

def extract_text_from_pptx(uploaded_file):
    try:
        prs = Presentation(uploaded_file); text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip(): text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"ERROR: {str(e)}"

def extract_text_from_docx(uploaded_file):
    try:
        doc = DocxDocument(uploaded_file); text = ""
        for para in doc.paragraphs:
            if para.text.strip(): text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip(): text += cell.text + "\n"
        return text.strip()
    except Exception as e:
        return f"ERROR: {str(e)}"

def extract_text_from_txt(uploaded_file):
    try: return uploaded_file.read().decode("utf-8", errors="ignore").strip()
    except Exception as e: return f"ERROR: {str(e)}"

def extract_text_from_file(uploaded_file):
    if uploaded_file is None: return None
    filename = uploaded_file.name.lower()
    if filename.endswith(".pdf"):  return extract_text_from_pdf(uploaded_file)
    if filename.endswith(".pptx"): return extract_text_from_pptx(uploaded_file)
    if filename.endswith(".docx"): return extract_text_from_docx(uploaded_file)
    if filename.endswith(".txt"):  return extract_text_from_txt(uploaded_file)
    return "ERROR: Unsupported file type."

def get_google_services():
    creds = None
    try:
        if "GOOGLE_TOKEN_JSON" in st.secrets:
            token_info = json.loads(st.secrets["GOOGLE_TOKEN_JSON"])
            creds = Credentials.from_authorized_user_info(token_info, SCOPES)
    except Exception: creds = None
    if creds is None and os.path.exists("token.json"):
        creds = Credentials.from_authorized_user_file("token.json", SCOPES)
    if not creds: st.error("No Google credentials found."); st.stop()
    if not creds.valid:
        if creds.expired and creds.refresh_token:
            creds.refresh(google.auth.transport.requests.Request())
            try: using_secrets = "GOOGLE_TOKEN_JSON" in st.secrets
            except Exception: using_secrets = False
            if not using_secrets:
                with open("token.json", "w") as f: f.write(creds.to_json())
        else: st.error("Google credentials invalid."); st.stop()
    return (build("forms","v1",credentials=creds), build("sheets","v4",credentials=creds), build("drive","v3",credentials=creds))

def create_google_form(questions, subject, difficulty, forms_service, drive_service):
    form_body = {"info":{"title":f"{subject} Quiz — {difficulty}","documentTitle":f"{subject} Quiz"}}
    form = forms_service.forms().create(body=form_body).execute()
    form_id = form["formId"]
    requests = []
    requests.append({"updateFormInfo":{"info":{"description":f"Subject: {subject}  |  Difficulty: {difficulty}  |  Total Marks: {len(questions)}\n\nInstructions: Select the correct answer for each question."},"updateMask":"description"}})
    requests.append({"createItem":{"item":{"title":"Student Full Name","description":"Enter your full name exactly as on your ID card.","questionItem":{"question":{"required":True,"textQuestion":{"paragraph":False}}}},"location":{"index":0}}})
    requests.append({"createItem":{"item":{"title":"Roll Number","description":"Enter your roll number (e.g. CS-24-01).","questionItem":{"question":{"required":True,"textQuestion":{"paragraph":False}}}},"location":{"index":1}}})
    for i, q in enumerate(questions):
        requests.append({"createItem":{"item":{"title":f"Q{i+1}. {q['question']}","questionItem":{"question":{"required":True,"choiceQuestion":{"type":"RADIO","options":[{"value":opt} for opt in q["options"]],"shuffle":True}}}},"location":{"index":i+2}}})
    forms_service.forms().batchUpdate(formId=form_id, body={"requests":requests}).execute()
    quiz_requests = [{"updateSettings":{"settings":{"quizSettings":{"isQuiz":True}},"updateMask":"quizSettings"}}]
    updated_form = forms_service.forms().get(formId=form_id).execute()
    items = updated_form.get("items",[])
    for i, q in enumerate(questions):
        item_index = i+2
        if item_index >= len(items): break
        item_id = items[item_index]["itemId"]
        question_id = items[item_index]["questionItem"]["question"]["questionId"]
        quiz_requests.append({"updateItem":{"item":{"itemId":item_id,"title":f"Q{i+1}. {q['question']}","questionItem":{"question":{"questionId":question_id,"required":True,"grading":{"pointValue":1,"correctAnswers":{"answers":[{"value":q["answer"]}]},"whenRight":{"text":f"Correct! {q['explanation']}"},"whenWrong":{"text":f"Wrong. Correct answer: {q['answer']}"}},"choiceQuestion":{"type":"RADIO","options":[{"value":opt} for opt in q["options"]],"shuffle":True}}}},"location":{"index":item_index},"updateMask":"title,questionItem"}})
    forms_service.forms().batchUpdate(formId=form_id, body={"requests":quiz_requests}).execute()
    drive_service.permissions().create(fileId=form_id, body={"type":"anyone","role":"reader"}).execute()
    return (form_id,
            f"https://docs.google.com/forms/d/{form_id}/viewform",
            f"https://docs.google.com/forms/d/{form_id}/edit",
            f"https://docs.google.com/forms/d/{form_id}/edit#responses")


# ══ PAGE ROUTER ══
page = st.session_state.page

# ─── SPLASH PAGE ─────────────────────────────────────────────────────────────
if page == "splash":

    st.markdown("""
    <style>
        [data-testid="stSidebar"]                { display: none !important; }
        [data-testid="stSidebarCollapsedControl"] { display: none !important; }
        [data-testid="stSidebarCollapseButton"]   { display: none !important; }
        .main .block-container { max-width: 100% !important; padding: 0 !important; }
    </style>
    """, unsafe_allow_html=True)

    st.html("""
<style>
.splash-root {
    min-height: 100vh;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    padding: 2.5rem 1rem 3rem 1rem;
    text-align: center;
    position: relative;
    overflow: hidden;
}
.splash-root::before {
    content: "";
    position: fixed;
    inset: 0;
    background-image:
        linear-gradient(rgba(0,229,255,0.04) 1px, transparent 1px),
        linear-gradient(90deg, rgba(0,229,255,0.04) 1px, transparent 1px);
    background-size: 40px 40px;
    pointer-events: none;
    animation: splash-grid-pulse 4s ease-in-out infinite;
    z-index: 0;
}
.splash-root::after {
    content: "";
    position: fixed;
    left: 0; width: 100%; height: 2px;
    background: linear-gradient(90deg, transparent 0%, rgba(0,229,255,0.55) 40%, rgba(124,58,255,0.45) 60%, transparent 100%);
    animation: splash-scan 7s ease-in-out infinite;
    pointer-events: none;
    z-index: 1;
}
.sp-corner { position: fixed; width: 70px; height: 70px; pointer-events: none; z-index: 2; }
.sp-corner.tl { top: 18px; left: 18px; border-top: 2px solid rgba(0,229,255,0.7); border-left: 2px solid rgba(0,229,255,0.7); animation: splash-corner-blink 3s ease-in-out infinite 0s; }
.sp-corner.tr { top: 18px; right: 18px; border-top: 2px solid rgba(0,229,255,0.7); border-right: 2px solid rgba(0,229,255,0.7); animation: splash-corner-blink 3s ease-in-out infinite 0.75s; }
.sp-corner.bl { bottom: 18px; left: 18px; border-bottom: 2px solid rgba(0,229,255,0.7); border-left: 2px solid rgba(0,229,255,0.7); animation: splash-corner-blink 3s ease-in-out infinite 1.5s; }
.sp-corner.br { bottom: 18px; right: 18px; border-bottom: 2px solid rgba(0,229,255,0.7); border-right: 2px solid rgba(0,229,255,0.7); animation: splash-corner-blink 3s ease-in-out infinite 2.25s; }
.sp-streak {
    position: fixed; top: 38%; height: 1px; width: 30%;
    background: linear-gradient(90deg, transparent, rgba(0,229,255,0.5), rgba(124,58,255,0.3), transparent);
    animation: splash-streak 5s ease-in-out infinite;
    pointer-events: none; z-index: 2;
}
.sp-streak.s2 { top: 62%; animation-delay: 2.5s; animation-duration: 6s; }
.splash-inner { position: relative; z-index: 3; width: 100%; max-width: 900px; margin: 0 auto; }
.sp-orb-wrap {
    position: relative;
    width: 160px; height: 160px;
    margin: 0 auto 2.4rem auto;
    display: flex; align-items: center; justify-content: center;
    animation: splash-float 5s ease-in-out infinite;
}
.sp-ring { position: absolute; inset: 0; border-radius: 50%; border-style: dashed; }
.sp-ring-1 { border: 1.5px dashed rgba(0,229,255,0.55); animation: splash-ring-1 6s linear infinite; }
.sp-ring-2 { inset: -16px; border: 1px dashed rgba(124,58,255,0.4); animation: splash-ring-2 10s linear infinite; }
.sp-ring-3 { inset: -32px; border: 1px dotted rgba(0,255,136,0.2); animation: splash-ring-3 14s linear infinite; }
.sp-hex-core {
    width: 110px; height: 110px;
    border-radius: 18px;
    background: linear-gradient(135deg, #0C1A1E 0%, #0F2128 50%, #142830 100%);
    border: 2px solid rgba(0,229,255,0.5);
    display: flex; align-items: center; justify-content: center;
    font-size: 3.6rem;
    box-shadow: 0 0 0 1px rgba(0,229,255,0.15), 0 0 30px rgba(0,229,255,0.25), 0 0 60px rgba(0,229,255,0.1), inset 0 0 20px rgba(0,229,255,0.08);
    animation: splash-hex-spin 12s ease-in-out infinite, avatar-pulse 3s ease-in-out infinite;
    position: relative; z-index: 1;
}
.sp-particle {
    position: absolute; width: 5px; height: 5px; border-radius: 50%;
    background: var(--cyan, #00E5FF);
    animation: splash-particle 3s ease-in-out infinite;
}
.sp-particle:nth-child(1) { top: 10%; left: 20%; animation-delay: 0s;   background: #00E5FF; }
.sp-particle:nth-child(2) { top: 20%; right: 15%; animation-delay: 0.8s; background: #7C3AFF; }
.sp-particle:nth-child(3) { bottom: 15%; left: 25%; animation-delay: 1.6s;background: #00FF88; }
.sp-particle:nth-child(4) { bottom: 20%; right: 20%; animation-delay: 2.4s;background: #00E5FF; }
.sp-title-badge-wrap {
    position: relative;
    display: inline-flex; align-items: center; justify-content: center;
    margin-bottom: 1.8rem;
    padding: 3px;
    border-radius: 10px;
    background: linear-gradient(90deg, #00E5FF, #7C3AFF, #00FF88, #00E5FF);
    background-size: 300% 100%;
    animation: splash-border-run 3s linear infinite, splash-float 6s ease-in-out infinite 1s;
}
.sp-title-badge {
    position: relative;
    background: #0C1A1E;
    border-radius: 8px;
    padding: 1.4rem 2.8rem;
    overflow: hidden;
}
.sp-title-badge::before {
    content: "";
    position: absolute; top: 0; left: 0; right: 0; bottom: 0;
    background: linear-gradient(110deg, transparent 25%, rgba(0,229,255,0.14) 50%, transparent 75%);
    background-size: 200% 100%;
    animation: badge-shimmer 2.5s linear infinite;
    border-radius: 8px;
    pointer-events: none;
}
.sp-badge-line1 {
    font-family: 'Share Tech Mono', monospace;
    font-size: clamp(1.5rem, 5vw, 3.2rem);
    letter-spacing: 6px;
    text-transform: uppercase;
    display: block;
    margin-bottom: 0.2rem;
    animation: splash-title-reveal 0.9s ease-out backwards;
}
.sp-badge-word-ai    { color: #E8F8FF; }
.sp-badge-word-quiz  { color: #00E5FF; text-shadow: 0 0 20px rgba(0,229,255,0.7), 0 0 40px rgba(0,229,255,0.3); }
.sp-badge-word-gen   { color: #7C3AFF; text-shadow: 0 0 20px rgba(124,58,255,0.6); }
.sp-badge-line2 {
    font-family: 'Share Tech Mono', monospace;
    font-size: clamp(1.1rem, 3.5vw, 2.2rem);
    letter-spacing: 5px;
    text-transform: uppercase;
    display: block;
    animation: splash-title-reveal 0.9s ease-out 0.15s backwards;
}
.sp-badge-word-and     { color: #C8F0F7; font-size: 0.75em; margin: 0 8px; }
.sp-badge-word-checker { color: #00FF88; text-shadow: 0 0 18px rgba(0,255,136,0.6); }
.sp-badge-bar {
    display: block; height: 2px; border-radius: 1px; margin: 0.9rem auto 0 auto;
    background: linear-gradient(90deg, transparent, #00E5FF, #7C3AFF, #00FF88, transparent);
    animation: splash-bar-slide 0.8s ease-out 0.5s backwards, border-rotate 3s linear infinite 1.5s;
    background-size: 200% 100%;
}
.sp-version-chip {
    display: inline-flex; align-items: center; gap: 10px;
    border: 1px solid rgba(0,229,255,0.3);
    border-radius: 3px; padding: 0.35rem 1.2rem;
    font-family: 'Share Tech Mono', monospace; font-size: 0.65rem;
    color: rgba(0,229,255,0.8); letter-spacing: 3px; text-transform: uppercase;
    margin-bottom: 2.2rem;
    animation: splash-sub-reveal 0.7s ease-out 0.6s backwards, pulse-cyan 5s ease-in-out infinite 2s;
}
.sp-vc-dot {
    width: 6px; height: 6px; border-radius: 50%;
    background: #00FF88; animation: status-ping 1.8s ease-in-out infinite;
    flex-shrink: 0;
}
.sp-creator-wrap {
    position: relative;
    display: inline-flex; align-items: center; justify-content: center;
    padding: 2px;
    border-radius: 8px;
    background: linear-gradient(135deg, rgba(0,229,255,0.4), rgba(124,58,255,0.4), rgba(0,255,136,0.3));
    background-size: 200% 200%;
    animation: splash-creator-in 0.8s ease-out 0.8s backwards, border-rotate 5s linear infinite 1.8s;
    margin-bottom: 2.2rem;
}
.sp-creator-badge {
    background: #0C1A1E;
    border-radius: 6px;
    padding: 1rem 2rem;
    display: flex; align-items: center; gap: 1.2rem;
    position: relative; overflow: hidden;
}
.sp-creator-badge::before {
    content: "";
    position: absolute; top: 0; left: 0; right: 0; bottom: 0;
    background: linear-gradient(120deg, transparent, rgba(124,58,255,0.08), transparent);
    animation: badge-shimmer 4s linear infinite;
    pointer-events: none;
}
.sp-av-wrap {
    position: relative; width: 56px; height: 56px; min-width: 56px;
    display: flex; align-items: center; justify-content: center;
}
.sp-av-wrap::before {
    content: ""; position: absolute;
    inset: -8px; border-radius: 10px;
    border: 1px dashed rgba(0,229,255,0.4);
    animation: ring-spin 7s linear infinite;
}
.sp-av-wrap::after {
    content: ""; position: absolute;
    inset: -15px; border-radius: 12px;
    border: 1px dotted rgba(124,58,255,0.3);
    animation: ring-spin-reverse 11s linear infinite;
}
.sp-avatar {
    width: 56px; height: 56px; border-radius: 6px;
    background: linear-gradient(135deg, #7C3AFF 0%, #3B5FE0 60%, #00E5FF 100%);
    display: flex; align-items: center; justify-content: center;
    font-family: 'Share Tech Mono', monospace; font-size: 1.2rem; font-weight: 700; color: #fff;
    border: 1px solid rgba(124,58,255,0.6);
    box-shadow: 0 0 20px rgba(124,58,255,0.4), 0 0 40px rgba(0,229,255,0.15);
    animation: avatar-pulse 2.8s ease-in-out infinite;
    position: relative; z-index: 1;
}
.sp-creator-text { text-align: left; position: relative; z-index: 1; }
.sp-creator-name {
    font-family: 'Share Tech Mono', monospace;
    font-size: clamp(1rem, 2.5vw, 1.25rem);
    color: #E8F8FF; letter-spacing: 2px;
    margin-bottom: 3px;
    animation: glow-pulse 4s ease-in-out infinite;
}
.sp-creator-meta {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.7rem; color: #3D6B78;
    letter-spacing: 1.5px; text-transform: uppercase;
    margin-bottom: 0.5rem;
}
.sp-creator-tags { display: flex; gap: 0.45rem; flex-wrap: wrap; }
.sp-ctag {
    font-family: 'JetBrains Mono', monospace; font-size: 0.6rem;
    padding: 3px 10px; border: 1px solid; border-radius: 2px;
    letter-spacing: 1px; text-transform: uppercase; font-weight: 600;
}
.sp-ctag-c { border-color: rgba(0,229,255,0.4);   color: #009AB0; animation: tag-glow-cycle 3s ease-in-out infinite 0s; }
.sp-ctag-g { border-color: rgba(0,255,136,0.35);  color: #00FF88; animation: tag-glow-cycle 3s ease-in-out infinite 1s; }
.sp-ctag-p { border-color: rgba(124,58,255,0.4);  color: #7C3AFF; animation: tag-glow-cycle 3s ease-in-out infinite 2s; }
.sp-ctag-a { border-color: rgba(255,203,61,0.4);  color: #FFCB3D; animation: tag-glow-cycle 3s ease-in-out infinite 0.5s; }
.sp-intro {
    font-family: 'JetBrains Mono', monospace;
    font-size: clamp(0.75rem, 1.5vw, 0.88rem);
    color: #3D6B78; letter-spacing: 0.8px; line-height: 1.85;
    max-width: 620px; margin: 0 auto 2.6rem auto;
    animation: splash-sub-reveal 0.7s ease-out 1s backwards;
}
.sp-intro strong { color: #00E5FF; }
.sp-btn-hint {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.62rem; color: #1D3A42; letter-spacing: 3px;
    text-transform: uppercase; margin-top: 0.6rem;
    animation: dot-blink 2s ease-in-out infinite;
}
</style>

<div class="splash-root">
    <div class="sp-corner tl"></div>
    <div class="sp-corner tr"></div>
    <div class="sp-corner bl"></div>
    <div class="sp-corner br"></div>
    <div class="sp-streak"></div>
    <div class="sp-streak s2"></div>
    <div class="splash-inner">
        <div class="sp-orb-wrap">
            <div class="sp-ring sp-ring-3"></div>
            <div class="sp-ring sp-ring-2"></div>
            <div class="sp-ring sp-ring-1"></div>
            <div class="sp-particle"></div>
            <div class="sp-particle"></div>
            <div class="sp-particle"></div>
            <div class="sp-particle"></div>
            <div class="sp-hex-core">⬡</div>
        </div>
        <div class="sp-title-badge-wrap">
            <div class="sp-title-badge">
                <span class="sp-badge-line1">
                    <span class="sp-badge-word-ai">AI&nbsp;</span>
                    <span class="sp-badge-word-quiz">QUIZ&nbsp;</span>
                    <span class="sp-badge-word-gen">GENERATOR</span>
                </span>
                <span class="sp-badge-line2">
                    <span class="sp-badge-word-and">&amp;</span>
                    <span class="sp-badge-word-checker">CHECKER</span>
                </span>
                <span class="sp-badge-bar" style="width:120px;"></span>
            </div>
        </div>
        <div class="sp-version-chip">
            <span class="sp-vc-dot"></span>
            QUEST_SYSTEM &nbsp;·&nbsp; v2.0 &nbsp;·&nbsp; GROQ LLaMA 3.1 · 70B
            <span class="sp-vc-dot"></span>
        </div>
        <div style="display:flex;justify-content:center;margin-bottom:2rem;">
            <div class="sp-creator-wrap">
                <div class="sp-creator-badge">
                    <div class="sp-av-wrap">
                        <div class="sp-avatar">AR</div>
                    </div>
                    <div class="sp-creator-text">
                        <div class="sp-creator-name">Abdul Rehman Raja</div>
                        <div class="sp-creator-meta">BSCS · FAST-NUCES Islamabad</div>
                        <div class="sp-creator-tags">
                            <span class="sp-ctag sp-ctag-c">Groq API</span>
                            <span class="sp-ctag sp-ctag-p">LLaMA 3.1 70B</span>
                            <span class="sp-ctag sp-ctag-g">Google Forms</span>
                            <span class="sp-ctag sp-ctag-a">Streamlit</span>
                        </div>
                    </div>
                </div>
            </div>
        </div>
        <p class="sp-intro">
            Upload <strong>notes, PDFs, PPTX or DOCX</strong> — or just type a topic.
            QUEST generates up to <strong>50 multiple-choice questions</strong>, deploys them
            as a <strong>self-grading Google Form</strong>, then auto-calculates every
            student's score and exports a formatted <strong>Excel report</strong>.
        </p>
        <div class="sp-btn-hint">▼ &nbsp; click to launch &nbsp; ▼</div>
    </div>
</div>
""")

    col1, col2, col3 = st.columns([1.2, 1.6, 1.2])
    with col2:
        if st.button("⬡  GET STARTED", use_container_width=True, type="primary"):
            st.session_state.page = "home"
            st.rerun()

    st.markdown(
        "<p style='text-align:center;font-family:JetBrains Mono,monospace;"
        "font-size:0.6rem;color:#1D3A42;letter-spacing:2px;margin-top:0.4rem;'>"
        "QUEST // AI QUIZ SYSTEM — FAST-NUCES 2025</p>",
        unsafe_allow_html=True
    )


# ─── HOME PAGE ───────────────────────────────────────────────────────────────
elif page == "home":
    with st.sidebar:
        st.markdown('<div class="sidebar-logo">⬡ QUEST</div>', unsafe_allow_html=True)
        st.markdown('<div class="sidebar-sub">AI Quiz System v2.0</div>', unsafe_allow_html=True)
        st.markdown('<div class="sidebar-status"><span class="status-dot"></span>SYSTEM_ONLINE</div>', unsafe_allow_html=True)
        st.markdown('<hr class="sidebar-divider">', unsafe_allow_html=True)

        render_quest_tracker(st.session_state.page)

        nav_items = [
            ("home",  "⬡", "[ HOME ]",     ""),
            ("step1", "01", "CREATE_QUIZ",  "AI"),
            ("step2", "02", "PUBLISH_FORM", "FORMS"),
            ("step3", "03", "CHECK_RESULTS","DATA"),
        ]

        for key, icon, label, badge in nav_items:
            is_active = (st.session_state.page == key)
            badge_part = f"  ·  {badge}" if badge else ""
            btn_label = f"{icon}  {label}{badge_part}"
            active_class = "nav-active" if is_active else "nav-inactive"
            st.markdown(f'<div class="{active_class}">', unsafe_allow_html=True)
            if st.button(btn_label, key=f"nav_{key}", use_container_width=True):
                st.session_state.page = key
                st.rerun()
            st.markdown('</div>', unsafe_allow_html=True)

        st.markdown('<hr class="sidebar-divider">', unsafe_allow_html=True)

        if "questions" in st.session_state:
            st.html(f"""<div style="background:rgba(0,255,136,0.06);border:1px solid rgba(0,255,136,0.2);
border-left:2px solid #00FF88;border-radius:3px;padding:0.5rem 0.8rem;margin-bottom:0.4rem;
font-family:'JetBrains Mono',monospace;font-size:0.65rem;color:#00FF88;letter-spacing:1px;">
▶ {len(st.session_state.questions)} QUESTIONS_READY
</div>""")
        if st.session_state.get("form_url"):
            st.html("""<div style="background:rgba(0,229,255,0.06);border:1px solid rgba(0,229,255,0.2);
border-left:2px solid #00E5FF;border-radius:3px;padding:0.5rem 0.8rem;margin-bottom:0.4rem;
font-family:'JetBrains Mono',monospace;font-size:0.65rem;color:#00E5FF;letter-spacing:1px;">
▶ FORM_DEPLOYED
</div>""")
        if st.session_state.get("results"):
            st.html(f"""<div style="background:rgba(124,58,255,0.08);border:1px solid rgba(124,58,255,0.25);
border-left:2px solid #7C3AFF;border-radius:3px;padding:0.5rem 0.8rem;margin-bottom:0.4rem;
font-family:'JetBrains Mono',monospace;font-size:0.65rem;color:#A78BFA;letter-spacing:1px;">
▶ {len(st.session_state.results)} RESPONSES_GRADED
</div>""")

        st.markdown("""<div class="sidebar-creator">
    DEVELOPED BY<br>
    <span>Abdul Rehman Raja</span><br>
    <span>FAST-NUCES · Groq LLaMA 3.1 70B</span>
</div>""", unsafe_allow_html=True)

    st.html("""<div class="hero-wrapper">
    <div class="hud-corner tl"></div>
    <div class="hud-corner tr"></div>
    <div class="hud-corner bl"></div>
    <div class="hud-corner br"></div>
    <div class="master-badge-wrap">
        <div class="master-badge">
            <span class="master-badge-icon">⬡</span>
            <span class="master-badge-text-1">AI</span>
            <span class="master-badge-text-2">QUIZ</span>
            <span class="master-badge-text-3">GENERATOR</span>
            <span class="master-badge-amp">&amp;</span>
            <span class="master-badge-text-2">CHECKER</span>
        </div>
    </div>
    <div class="system-badge">
        <span class="badge-dot"></span>
        QUEST_SYSTEM · AI_QUIZ_GENERATOR
        <span class="badge-dot"></span>
    </div>
    <div class="hero-title">Generate. Deploy.</div>
    <span class="hero-title-accent">Grade Everything.<span class="cursor-blink"></span></span>
    <p class="hero-desc">
        Feed it your <strong>notes, a PDF, or just a topic</strong>. The AI writes
        up to <strong>50 multiple-choice questions</strong>, publishes them as a
        <strong>self-grading Google Form</strong>, then auto-calculates every
        student's score and exports a formatted Excel report.
    </p>
    <div class="flow-strip">
        <div class="flow-card">
            <span class="flow-num">STEP_01</span>
            <span class="flow-icon">⬡</span>
            <div class="flow-title">AI Generates</div>
            <div class="flow-desc">Paste notes or upload PDF/PPTX/DOCX. Choose up to 50 questions and difficulty.</div>
        </div>
        <div class="flow-sep"></div>
        <div class="flow-card">
            <span class="flow-num">STEP_02</span>
            <span class="flow-icon">◈</span>
            <div class="flow-title">You Deploy</div>
            <div class="flow-desc">One click creates a live, self-grading Google Form. Share the link.</div>
        </div>
        <div class="flow-sep"></div>
        <div class="flow-card">
            <span class="flow-num">STEP_03</span>
            <span class="flow-icon">◎</span>
            <div class="flow-title">Auto Grades</div>
            <div class="flow-desc">Fetch all responses, view scores, download Excel results sheet.</div>
        </div>
    </div>
    <div class="stat-strip">
        <div class="stat-pill"><span class="stat-num">5–50</span><span class="stat-label">QUESTIONS</span></div>
        <div class="stat-pill"><span class="stat-num">4</span><span class="stat-label">FILE_TYPES</span></div>
        <div class="stat-pill"><span class="stat-num">1-CLICK</span><span class="stat-label">PUBLISH</span></div>
        <div class="stat-pill"><span class="stat-num">AUTO</span><span class="stat-label">GRADING</span></div>
    </div>
    <div class="creator-card">
        <div class="creator-avatar-wrap">
            <div class="creator-avatar">AR</div>
        </div>
        <div>
            <div class="creator-name">Abdul Rehman Raja</div>
            <div class="creator-role">BSCS · FAST-NUCES Islamabad</div>
            <div class="creator-tags">
                <span class="ctag ctag-cyan">Groq API</span>
                <span class="ctag ctag-green">Google Forms</span>
                <span class="ctag ctag-purp">LLaMA 3.1 70B</span>
            </div>
        </div>
    </div>
    <div class="hero-stack">
        <span class="stack-pill">Groq LLaMA 3.1 70B</span>
        <span class="stack-pill">PDF · PPTX · DOCX · TXT</span>
        <span class="stack-pill">Google Forms API</span>
        <span class="stack-pill">Excel Export</span>
        <span class="stack-pill">Streamlit</span>
    </div>
</div>""")

    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        if st.button("[ INITIALIZE_QUIZ ]", use_container_width=True, type="primary"):
            st.session_state.page = "step1"
            st.rerun()


elif page in ("step1", "step2", "step3"):
    with st.sidebar:
        st.markdown('<div class="sidebar-logo">⬡ QUEST</div>', unsafe_allow_html=True)
        st.markdown('<div class="sidebar-sub">AI Quiz System v2.0</div>', unsafe_allow_html=True)
        st.markdown('<div class="sidebar-status"><span class="status-dot"></span>SYSTEM_ONLINE</div>', unsafe_allow_html=True)
        st.markdown('<hr class="sidebar-divider">', unsafe_allow_html=True)

        render_quest_tracker(st.session_state.page)

        nav_items = [
            ("home",  "⬡", "[ HOME ]",     ""),
            ("step1", "01", "CREATE_QUIZ",  "AI"),
            ("step2", "02", "PUBLISH_FORM", "FORMS"),
            ("step3", "03", "CHECK_RESULTS","DATA"),
        ]

        for key, icon, label, badge in nav_items:
            is_active = (st.session_state.page == key)
            badge_part = f"  ·  {badge}" if badge else ""
            btn_label = f"{icon}  {label}{badge_part}"
            active_class = "nav-active" if is_active else "nav-inactive"
            st.markdown(f'<div class="{active_class}">', unsafe_allow_html=True)
            if st.button(btn_label, key=f"nav_{key}", use_container_width=True):
                st.session_state.page = key
                st.rerun()
            st.markdown('</div>', unsafe_allow_html=True)

        st.markdown('<hr class="sidebar-divider">', unsafe_allow_html=True)

        if "questions" in st.session_state:
            st.html(f"""<div style="background:rgba(0,255,136,0.06);border:1px solid rgba(0,255,136,0.2);
border-left:2px solid #00FF88;border-radius:3px;padding:0.5rem 0.8rem;margin-bottom:0.4rem;
font-family:'JetBrains Mono',monospace;font-size:0.65rem;color:#00FF88;letter-spacing:1px;">
▶ {len(st.session_state.questions)} QUESTIONS_READY
</div>""")
        if st.session_state.get("form_url"):
            st.html("""<div style="background:rgba(0,229,255,0.06);border:1px solid rgba(0,229,255,0.2);
border-left:2px solid #00E5FF;border-radius:3px;padding:0.5rem 0.8rem;margin-bottom:0.4rem;
font-family:'JetBrains Mono',monospace;font-size:0.65rem;color:#00E5FF;letter-spacing:1px;">
▶ FORM_DEPLOYED
</div>""")
        if st.session_state.get("results"):
            st.html(f"""<div style="background:rgba(124,58,255,0.08);border:1px solid rgba(124,58,255,0.25);
border-left:2px solid #7C3AFF;border-radius:3px;padding:0.5rem 0.8rem;margin-bottom:0.4rem;
font-family:'JetBrains Mono',monospace;font-size:0.65rem;color:#A78BFA;letter-spacing:1px;">
▶ {len(st.session_state.results)} RESPONSES_GRADED
</div>""")

        st.markdown("""<div class="sidebar-creator">
    DEVELOPED BY<br>
    <span>Abdul Rehman Raja</span><br>
    <span>FAST-NUCES · Groq LLaMA 3.1 70B</span>
</div>""", unsafe_allow_html=True)

    # ─── STEP 1 — GENERATE ───────────────────────────────────────────────────────
    if page == "step1":
        render_journey_bar("step1")
        st.html("""<div class="page-header">
    <div class="page-title"><span>CREATE_QUIZ</span></div>
    <div class="page-subtitle">provide content source · configure parameters · execute generation</div>
</div>""")

        st.markdown('<div class="term-card">', unsafe_allow_html=True)
        st.markdown('<div class="sec-title">INPUT_METHOD</div>', unsafe_allow_html=True)
        input_method = st.radio(
            "Input method",
            ["PASTE_TEXT", "UPLOAD_FILE"],
            horizontal=True,
            label_visibility="collapsed"
        )
        st.markdown('</div>', unsafe_allow_html=True)

        extracted_content = ""

        if input_method == "PASTE_TEXT":
            st.markdown('<div class="term-card">', unsafe_allow_html=True)
            st.markdown('<div class="sec-title">TEXT_INPUT</div>', unsafe_allow_html=True)
            notes = st.text_area(
                "notes_input",
                placeholder="paste lecture notes, topic summary, or just a topic name...\ne.g. Newton's laws, Python OOP, World War II timeline",
                height=220,
                label_visibility="collapsed"
            )
            extracted_content = notes
            if notes.strip():
                wc = len(notes.split())
                st.caption(f"{wc:,} tokens loaded")
            st.markdown('</div>', unsafe_allow_html=True)
        else:
            st.markdown('<div class="term-card">', unsafe_allow_html=True)
            st.markdown('<div class="sec-title">FILE_UPLOAD</div>', unsafe_allow_html=True)
            uploaded_file = st.file_uploader(
                "upload",
                type=["pdf","pptx","docx","txt"],
                help="Supported: PDF, PowerPoint, Word, Plain Text",
                label_visibility="collapsed"
            )
            if uploaded_file is not None:
                file_size_kb = round(uploaded_file.size/1024, 1)
                st.html(f"""<div style="background:rgba(0,255,136,0.06);border:1px solid rgba(0,255,136,0.2);
border-left:2px solid #00FF88;border-radius:3px;padding:0.7rem 1rem;margin:0.5rem 0;
font-family:'JetBrains Mono',monospace;color:#00FF88;font-size:0.75rem;letter-spacing:1px;">
▶ FILE_LOADED: {uploaded_file.name} · {file_size_kb} KB
</div>""")
                with st.spinner("parsing file..."):
                    extracted_content = extract_text_from_file(uploaded_file)
                if extracted_content is None:
                    extracted_content = ""
                if extracted_content.startswith("ERROR:"):
                    st.error(f"{extracted_content}"); extracted_content = ""
                elif not extracted_content:
                    st.warning("no readable text found in file")
                else:
                    word_count = len(extracted_content.split())
                    c1, c2 = st.columns(2)
                    c1.metric("WORDS_FOUND", f"{word_count:,}")
                    c2.metric("CHARACTERS", f"{len(extracted_content):,}")
                    with st.expander("PREVIEW_EXTRACTED_TEXT"):
                        st.html(f"""<div style="background:var(--bg-raised, #0F2128);border:1px solid rgba(0,229,255,0.1);
border-radius:3px;padding:1rem;font-family:'JetBrains Mono',monospace;font-size:0.72rem;
color:#3D6B78;max-height:250px;overflow-y:auto;white-space:pre-wrap;letter-spacing:0.3px;">
{extracted_content[:1000]}{'...' if len(extracted_content)>1000 else ''}</div>""")
                    st.success("file parsed successfully")
            else:
                st.html("""<div class="dropzone">
<span class="dz-icon">⬡</span>
DROP_FILE_HERE or click to browse<br>
<span style="font-size:0.68rem;letter-spacing:2px;">PDF · PPTX · DOCX · TXT</span>
</div>""")
            st.markdown('</div>', unsafe_allow_html=True)

        st.markdown('<div class="term-card">', unsafe_allow_html=True)
        st.markdown('<div class="sec-title">QUIZ_PARAMETERS</div>', unsafe_allow_html=True)
        c1, c2, c3 = st.columns(3)
        with c1:
            # ── CHANGED: options now go up to 50; added 25, 30, 40, 50 ──
            num_q = st.selectbox("NUM_QUESTIONS", [5, 10, 15, 20, 25, 30, 40, 50], index=1)
        with c2:
            difficulty = st.selectbox("DIFFICULTY_LEVEL", ["Easy","Medium","Hard","Mixed"])
        with c3:
            subject = st.text_input("SUBJECT_NAME", placeholder="e.g. Physics, OOP, History")

        # Info banner for large question counts
        if num_q >= 30:
            st.html(f"""<div style="background:rgba(255,203,61,0.07);border:1px solid rgba(255,203,61,0.25);
border-left:2px solid #FFCB3D;border-radius:3px;padding:0.6rem 1rem;margin-top:0.6rem;
font-family:'JetBrains Mono',monospace;font-size:0.72rem;color:#FFCB3D;letter-spacing:0.8px;">
⚡ {num_q} questions selected — using LLaMA 3.1 70B for maximum accuracy.
Generation may take 15–30 seconds. Ensure your content is rich enough to cover {num_q} distinct concepts.
</div>""")

        st.markdown('</div>', unsafe_allow_html=True)

        col1, col2, col3 = st.columns([1,2,1])
        with col2:
            gen_btn = st.button("[ EXECUTE_GENERATION ]", use_container_width=True, type="primary")

        if gen_btn:
            if not extracted_content.strip():
                st.warning("ERROR: no content source provided")
            elif not subject.strip():
                st.warning("ERROR: subject name required")
            else:
                # Smart truncation — keeps full sentences, up to 12 000 chars
                content_for_ai = smart_truncate(extracted_content, max_chars=12000)

                progress_bar = st.progress(0, text="initializing model...")
                boot_msgs = [
                    "loading LLaMA 3.1 70B...",
                    "analysing content depth...",
                    "generating question pool...",
                    "selecting correct answers...",
                    "validating & appending explanations...",
                ]
                for i, msg in enumerate(boot_msgs):
                    progress_bar.progress(int((i+1) / (len(boot_msgs)+1) * 80), text=msg)
                    time.sleep(0.3)

                try:
                    # ── CHANGED: single clean call to the retry wrapper ──
                    questions = generate_questions_with_retry(
                        subject=subject,
                        difficulty=difficulty,
                        num_q=num_q,
                        content=content_for_ai,
                        max_retries=3,
                    )
                    progress_bar.progress(100, text="generation complete")
                    time.sleep(0.25)
                    progress_bar.empty()

                    st.session_state.questions  = questions
                    st.session_state.subject    = subject
                    st.session_state.difficulty = difficulty
                    st.session_state.form_id    = None
                    st.session_state.form_url   = None

                    st.html(f"""<div class="celebrate-banner">
    <div class="celebrate-title">GENERATION_COMPLETE</div>
    <div class="celebrate-sub">{len(questions)} questions compiled · model: LLaMA 3.1 70B · proceed to deployment</div>
</div>""")
                    st.balloons()

                except RuntimeError as e:
                    progress_bar.empty()
                    st.error(f"GENERATION_FAILED: {str(e)}")
                except Exception as e:
                    progress_bar.empty()
                    st.error(f"UNEXPECTED_ERROR: {str(e)}")

        if "questions" in st.session_state:
            questions = st.session_state.questions
            st.divider()
            st.markdown(
                f'<div class="sec-title">QUESTION_PREVIEW — {st.session_state.subject.upper()} '
                f'· {st.session_state.difficulty.upper()} · {len(questions)} QUESTIONS</div>',
                unsafe_allow_html=True
            )
            st.caption("highlighted option = correct answer  |  badges show concept & difficulty tag")

            for i, q in enumerate(questions):
                opts_html = ""
                for opt in q["options"]:
                    if opt == q["answer"]:
                        opts_html += f'<div class="q-correct">▶ {opt}</div>'
                    else:
                        opts_html += f'<div class="q-opt">◦ {opt}</div>'

                # ── CHANGED: display difficulty_tag + concept badges if present ──
                badges_html = ""
                if q.get("difficulty_tag"):
                    badges_html += f'<span class="ak-badge">{q["difficulty_tag"].upper()}</span>'
                if q.get("concept"):
                    badges_html += f'<span class="ak-badge" style="border-color:rgba(124,58,255,0.35);color:#7C3AFF;">{q["concept"]}</span>'

                st.html(f"""<div class="q-card" style="animation-delay:{min(i,8)*0.04}s;">
    <div class="q-title">
        <span class="q-num-badge">{i+1:02d}</span>
        {q['question']}
        {badges_html}
    </div>
    {opts_html}
    <div class="q-exp">{q['explanation']}</div>
</div>""")

            st.divider()
            col1, col2, col3 = st.columns([1,2,1])
            with col2:
                if st.button("[ PROCEED_TO_DEPLOYMENT ]", use_container_width=True, type="primary"):
                    st.session_state.page = "step2"
                    st.rerun()
        else:
            st.markdown("""<div class="empty-state">
<span class="empty-emoji">⬡</span>
AWAITING_INPUT<br>
<span style="font-size:0.72rem;">configure parameters above and execute generation</span>
</div>""", unsafe_allow_html=True)

    # ─── STEP 2 — PUBLISH ────────────────────────────────────────────────────────
    elif page == "step2":
        render_journey_bar("step2")
        st.html("""<div class="page-header">
    <div class="page-title"><span>PUBLISH_FORM</span></div>
    <div class="page-subtitle">deploy quiz to Google Forms · configure access · distribute link</div>
</div>""")

        if "questions" not in st.session_state:
            st.markdown("""<div class="empty-state">
<span class="empty-emoji">◈</span>
NO_QUIZ_FOUND<br>
<span style="font-size:0.72rem;">return to CREATE_QUIZ and generate questions first</span>
</div>""", unsafe_allow_html=True)
            col1,col2,col3 = st.columns([1,2,1])
            with col2:
                if st.button("[ BACK_TO_STEP_01 ]", use_container_width=True, type="primary"):
                    st.session_state.page = "step1"; st.rerun()
        else:
            st.markdown('<div class="term-card">', unsafe_allow_html=True)
            st.markdown('<div class="sec-title">QUIZ_MANIFEST</div>', unsafe_allow_html=True)
            c1,c2,c3 = st.columns(3)
            c1.metric("QUESTIONS", len(st.session_state.questions))
            c2.metric("TOTAL_MARKS", len(st.session_state.questions))
            c3.metric("DIFFICULTY", st.session_state.difficulty.upper())
            st.markdown('</div>', unsafe_allow_html=True)

            st.markdown('<div class="term-card">', unsafe_allow_html=True)
            st.markdown('<div class="sec-title">DEPLOY_TO_GOOGLE_FORMS</div>', unsafe_allow_html=True)
            st.info("First run: browser window opens for Google OAuth — automatic after that")
            col1,col2,col3 = st.columns([1,2,1])
            with col2:
                pub_btn = st.button("[ DEPLOY_FORM ]", use_container_width=True, type="primary")
            st.markdown('</div>', unsafe_allow_html=True)

            if pub_btn:
                progress_bar = st.progress(0, text="connecting to Google APIs...")
                try:
                    progress_bar.progress(25, text="authenticating...")
                    forms_svc, sheets_svc, drive_svc = get_google_services()
                    progress_bar.progress(55, text="building form structure...")
                    form_id, form_url, edit_url, responses_url = create_google_form(
                        st.session_state.questions, st.session_state.subject,
                        st.session_state.difficulty, forms_svc, drive_svc
                    )
                    progress_bar.progress(90, text="setting permissions...")
                    time.sleep(0.2)
                    st.session_state.form_id       = form_id
                    st.session_state.form_url      = form_url
                    st.session_state.edit_url      = edit_url
                    st.session_state.responses_url = responses_url
                    progress_bar.progress(100, text="deployed")
                    time.sleep(0.2)
                    progress_bar.empty()
                    st.html("""<div class="celebrate-banner">
    <div class="celebrate-title">FORM_DEPLOYED</div>
    <div class="celebrate-sub">copy the URL below and distribute to students</div>
</div>""")
                    st.balloons()
                except Exception as e:
                    progress_bar.empty()
                    st.error(f"DEPLOY_ERROR: {str(e)}")

            if st.session_state.get("form_url"):
                st.divider()
                st.markdown('<div class="sec-title">ACCESS_LINKS</div>', unsafe_allow_html=True)
                st.markdown('<div class="term-card">', unsafe_allow_html=True)
                st.code(st.session_state.form_url, language=None)
                c1,c2 = st.columns(2)
                with c1: st.link_button("[ STUDENT_VIEW ]", st.session_state.form_url, use_container_width=True)
                with c2: st.link_button("[ TEACHER_EDIT ]", st.session_state.edit_url, use_container_width=True)
                st.markdown('</div>', unsafe_allow_html=True)

                st.markdown('<div class="sec-title">ANSWER_KEY</div>', unsafe_allow_html=True)
                st.caption("classified — teacher access only")
                with st.expander("REVEAL_ANSWER_KEY", expanded=False):
                    for i, q in enumerate(st.session_state.questions):
                        st.html(f"""<div class="q-card" style="animation-delay:{min(i,8)*0.04}s;">
    <div class="q-title"><span class="q-num-badge">{i+1:02d}</span>{q['question']}</div>
    <div class="q-correct">▶ {q['answer']} <span class="ak-badge">CORRECT</span></div>
</div>""")

                st.divider()
                col1,col2,col3 = st.columns([1,2,1])
                with col2:
                    if st.button("[ FETCH_RESULTS ]", use_container_width=True, type="primary"):
                        st.session_state.page = "step3"; st.rerun()

    # ─── STEP 3 — RESULTS ────────────────────────────────────────────────────────
    elif page == "step3":
        render_journey_bar("step3")
        st.html("""<div class="page-header">
    <div class="page-title"><span>CHECK_RESULTS</span></div>
    <div class="page-subtitle">fetch responses · calculate grades · export report</div>
</div>""")

        if "questions" not in st.session_state:
            st.markdown("""<div class="empty-state">
<span class="empty-emoji">◎</span>
NO_QUIZ_FOUND<br>
<span style="font-size:0.72rem;">complete CREATE_QUIZ step first</span>
</div>""", unsafe_allow_html=True)
            col1,col2,col3 = st.columns([1,2,1])
            with col2:
                if st.button("[ BACK_TO_STEP_01 ]", use_container_width=True, type="primary"):
                    st.session_state.page = "step1"; st.rerun()
        elif not st.session_state.get("form_id"):
            st.markdown("""<div class="empty-state">
<span class="empty-emoji">◈</span>
FORM_NOT_DEPLOYED<br>
<span style="font-size:0.72rem;">complete PUBLISH_FORM step first</span>
</div>""", unsafe_allow_html=True)
            col1,col2,col3 = st.columns([1,2,1])
            with col2:
                if st.button("[ BACK_TO_STEP_02 ]", use_container_width=True, type="primary"):
                    st.session_state.page = "step2"; st.rerun()
        else:
            st.markdown('<div class="term-card">', unsafe_allow_html=True)
            st.markdown('<div class="sec-title">LIVE_MONITORING</div>', unsafe_allow_html=True)
            st.link_button("[ OPEN_RESPONSES_DASHBOARD ]", st.session_state.responses_url, use_container_width=True)
            st.markdown('</div>', unsafe_allow_html=True)

            col1,col2,col3 = st.columns([1,2,1])
            with col2:
                fetch_btn = st.button("[ FETCH_AND_GRADE ]", use_container_width=True, type="primary")

            if fetch_btn:
                progress_bar = st.progress(0, text="connecting to Google Forms...")
                try:
                    progress_bar.progress(20, text="authenticating...")
                    forms_svc, sheets_svc, drive_svc = get_google_services()
                    progress_bar.progress(45, text="downloading responses...")
                    resp = forms_svc.forms().responses().list(formId=st.session_state.form_id).execute()
                    responses = resp.get("responses",[])
                    if not responses:
                        progress_bar.empty()
                        st.warning("NO_RESPONSES_FOUND — share form link with students first")
                    else:
                        progress_bar.progress(70, text="running grading algorithm...")
                        questions = st.session_state.questions
                        total = len(questions)
                        form_data = forms_svc.forms().get(formId=st.session_state.form_id).execute()
                        items = form_data.get("items",[])
                        name_qid = items[0]["questionItem"]["question"]["questionId"]
                        roll_qid = items[1]["questionItem"]["question"]["questionId"]
                        mcq_qids = [items[i+2]["questionItem"]["question"]["questionId"]
                                    for i in range(len(questions)) if i+2 < len(items)]
                        results = []
                        for r in responses:
                            answers = r.get("answers",{})
                            name = answers.get(name_qid,{}).get("textAnswers",{}).get("answers",[{}])[0].get("value","Unknown")
                            roll = answers.get(roll_qid,{}).get("textAnswers",{}).get("answers",[{}])[0].get("value","N/A")
                            marks = sum(
                                1 for j, qid in enumerate(mcq_qids)
                                if answers.get(qid,{}).get("textAnswers",{}).get("answers",[{}])[0].get("value","") == questions[j]["answer"]
                            )
                            results.append({"name":name,"roll":roll,"marks":marks})
                        st.session_state.results = results
                        st.session_state.total   = total
                        progress_bar.progress(100, text="grading complete")
                        time.sleep(0.2)
                        progress_bar.empty()
                        st.html(f"""<div class="celebrate-banner">
    <div class="celebrate-title">GRADING_COMPLETE</div>
    <div class="celebrate-sub">{len(results)} responses processed · download report below</div>
</div>""")
                        st.balloons()
                except Exception as e:
                    progress_bar.empty()
                    st.error(f"FETCH_ERROR: {str(e)}")

            if st.session_state.get("results"):
                results = st.session_state.results
                total   = st.session_state.total
                st.divider()
                passed = sum(1 for r in results if (r["marks"]/total)*100 >= 50)
                avg    = round(sum(r["marks"] for r in results)/len(results), 1)

                st.markdown('<div class="term-card">', unsafe_allow_html=True)
                st.markdown('<div class="sec-title">CLASS_ANALYTICS</div>', unsafe_allow_html=True)
                c1,c2,c3,c4 = st.columns(4)
                c1.metric("TOTAL_STUDENTS", len(results))
                c2.metric("PASSED", f"{passed}/{len(results)}")
                c3.metric("FAILED", f"{len(results)-passed}/{len(results)}")
                c4.metric("CLASS_AVG", f"{avg}/{total}")
                st.markdown('</div>', unsafe_allow_html=True)

                st.markdown('<div class="sec-title">RESULT_MATRIX</div>', unsafe_allow_html=True)

                sort_choice = st.radio(
                    "Sort by",
                    ["SCORE_DESC", "NAME_ASC", "ROLL_NUM"],
                    horizontal=True,
                    label_visibility="collapsed"
                )
                if sort_choice == "SCORE_DESC":
                    sorted_results = sorted(results, key=lambda r: r["marks"], reverse=True)
                elif sort_choice == "NAME_ASC":
                    sorted_results = sorted(results, key=lambda r: r["name"].lower())
                else:
                    sorted_results = sorted(results, key=lambda r: r["roll"])

                for i, r in enumerate(sorted_results):
                    pct = round((r["marks"]/total)*100, 1)
                    grade = get_grade(pct)
                    sym = get_grade_emoji(pct)
                    status_cls = "pass" if pct >= 50 else "fail"
                    st.html(f"""<div class="result-row {status_cls}" style="animation-delay:{min(i,10)*0.03}s;">
    <div style="flex:0 0 24px;font-family:'Share Tech Mono',monospace;font-size:0.8rem;color:var(--ink-muted, #3D6B78);">{sym}</div>
    <div style="flex:2;font-family:'JetBrains Mono',monospace;font-weight:500;font-size:0.82rem;color:var(--ink-bright, #E8F8FF);">{r['name']}</div>
    <div style="flex:1;font-family:'JetBrains Mono',monospace;font-size:0.72rem;color:var(--ink-muted, #3D6B78);letter-spacing:0.5px;">{r['roll']}</div>
    <div style="flex:1;font-family:'Share Tech Mono',monospace;font-size:0.82rem;color:var(--ink, #C8F0F7);">{r['marks']}/{total}</div>
    <div style="flex:1;font-family:'JetBrains Mono',monospace;font-size:0.75rem;color:var(--ink-muted, #3D6B78);">{pct}%</div>
    <div style="flex:0 0 52px;text-align:right;"><span class="grade-pill {status_cls}">{grade}</span></div>
</div>""")

                st.divider()
                col1,col2,col3 = st.columns([1,2,1])
                with col2:
                    excel_bytes = make_excel(results, st.session_state.subject, total)
                    st.download_button(
                        label="[ DOWNLOAD_EXCEL_REPORT ]",
                        data=excel_bytes,
                        file_name=f"{st.session_state.subject}_results.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True
                    )
